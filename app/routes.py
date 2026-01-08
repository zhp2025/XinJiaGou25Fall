from flask import Blueprint, render_template, request, jsonify, redirect, url_for, flash, session
from flask_login import login_user, logout_user, login_required, current_user
from app.mock_data import (
    MOCK_ARTICLES, MOCK_TOOLS, MOCK_CASES, MOCK_TERMS, 
    MOCK_RESOURCES, MOCK_ETHICS_TOPICS, MOCK_FORUM_POSTS,
    MOCK_LEARNING_PATHS,
    get_user_by_username, get_user_by_id, get_all_terms_from_db
)
from app import MockUser
from datetime import datetime, timedelta

# 全局反馈列表（临时存储，实际应使用数据库）
feedback_list = []

# 注意：FORUM_LIKES已废弃，统一使用MOCK_LIKES

def record_visit():
    if current_user.is_authenticated:
        from app import record_visit as record_visit_func
        record_visit_func(current_user.id)

# 主蓝图
main_bp = Blueprint('main', __name__)

# 全局登录要求 - 除了登录和注册页面，所有页面都需要登录
@main_bp.before_request
def require_login():
    """要求所有页面都需要登录（除了登录和注册页面）"""
    from flask_login import current_user
    # 排除登录和注册相关的路由和静态文件
    excluded_paths = ['/auth/login', '/auth/register', '/auth/forgot-password']
    if request.path.startswith('/static/'):
        return
    if request.path not in excluded_paths and not current_user.is_authenticated:
        from flask import redirect, url_for, session
        # 保存原始URL以便登录后重定向
        session['next_url'] = request.url
        return redirect(url_for('auth.login'))

@main_bp.route('/')
def index():
    """主页"""
    record_visit()
    
    # 获取热门科普文章（使用标记系统规范化分类）
    from app.category_utils import filter_items_by_category
    from app.tagging_system import normalize_article_category
    
    # 规范化"热门科普"为"科普文章"
    popular_articles = filter_items_by_category(
        [a for a in MOCK_ARTICLES if a.get('is_featured', False)],
        '科普文章',
        'article'
    )
    popular_articles = sorted(popular_articles, key=lambda x: x['views'], reverse=True)[:6]
    
    # 获取最新资讯（规范化"最新资讯"为"资讯动态"）
    latest_news = filter_items_by_category(MOCK_ARTICLES, '资讯动态', 'article')
    latest_news = sorted(latest_news, key=lambda x: x['created_at'], reverse=True)[:6]
    
    # 获取热门大模型（从工具中筛选，使用LLM分类）
    popular_models = filter_items_by_category(MOCK_TOOLS, 'LLM', 'tool')
    popular_models = sorted(popular_models, key=lambda x: x['rating'], reverse=True)[:4]
    
    # 获取热门帖子（用于社区排行榜）
    hot_posts = sorted(MOCK_FORUM_POSTS, key=lambda x: (x['likes'], x['views']), reverse=True)[:5]
    # 为帖子添加作者信息
    for post in hot_posts:
        user = get_user_by_id(post['user_id'])
        if user:
            display_name = user.get('nickname', user['username'])
            post['author'] = {'username': display_name}
        else:
            post['author'] = {'username': '未知用户'}
        post['comments_count'] = post.get('comments_count', 0)
    
    # 获取资源中心数据（前3个资源）
    resources_preview = sorted(MOCK_RESOURCES, key=lambda x: x['created_at'], reverse=True)[:3]
    
    # 获取应用场景数据（前3个案例）
    cases_preview = sorted(MOCK_CASES, key=lambda x: x['created_at'], reverse=True)[:3]
    
    # 获取伦理与未来数据（前3个专题）
    ethics_preview = sorted(MOCK_ETHICS_TOPICS, key=lambda x: x['created_at'], reverse=True)[:3]
    
    return render_template('index.html', 
                         popular_articles=popular_articles,
                         latest_news=latest_news,
                         popular_models=popular_models,
                         hot_posts=hot_posts,
                         resources_preview=resources_preview,
                         cases_preview=cases_preview,
                         ethics_preview=ethics_preview)


@main_bp.route('/ai-basics')
def ai_basics():
    """AI基础模块"""
    record_visit()
    
    # 从数据库获取所有术语（优先从数据库）
    all_terms = get_all_terms_from_db()
    
    # 获取核心概念（术语）- 优先选择核心概念类别，如果不够8个，再从其他类别补充
    from app.tagging_system import normalize_term_category
    
    core_concepts = [t for t in all_terms if normalize_term_category(t.get('category', '')) == '核心概念']
    # 如果核心概念不够8个，从其他类别补充
    if len(core_concepts) < 8:
        other_categories = ['LLM', 'Transformer', '扩散模型']
        other_concepts = [
            t for t in all_terms 
            if normalize_term_category(t.get('category', '')) in other_categories and t not in core_concepts
        ]
        core_concepts.extend(other_concepts[:8-len(core_concepts)])
    # 只取前8个
    core_concepts = core_concepts[:8]
    
    # 将所有概念都归类为文字概念（不再区分图片、视频等）
    text_concepts = [c for c in core_concepts if c.get('definition')]
    # 不再使用图片概念、视频概念和知识图谱概念的分类
    image_concepts = []
    video_concepts = []
    knowledge_graph_concepts = []
    
    return render_template('ai_basics.html', 
                         text_concepts=text_concepts,
                         image_concepts=image_concepts,
                         video_concepts=video_concepts,
                         knowledge_graph_concepts=knowledge_graph_concepts,
                         learning_paths=MOCK_LEARNING_PATHS)


@main_bp.route('/ai-lab')
def ai_lab():
    """AI实验室模块"""
    record_visit()
    return render_template('ai_lab.html')


@main_bp.route('/applications')
def applications():
    """应用场景模块"""
    record_visit()
    
    # 分页参数
    page = request.args.get('page', 1, type=int)
    per_page = 9
    
    # 获取所有案例
    all_cases = sorted(MOCK_CASES, key=lambda x: x['created_at'], reverse=True)
    
    # 分页处理
    total = len(all_cases)
    start = (page - 1) * per_page
    end = start + per_page
    cases = all_cases[start:end]
    
    # 创建分页对象
    class Pagination:
        def __init__(self, items, page, per_page, total):
            self.items = items
            self.page = page
            self.per_page = per_page
            self.total = total
            self.pages = (total + per_page - 1) // per_page if total > 0 else 1
            self.has_prev = page > 1
            self.has_next = page < self.pages
            self.prev_num = page - 1 if self.has_prev else None
            self.next_num = page + 1 if self.has_next else None
        
        def iter_pages(self, left_edge=1, right_edge=1, left_current=2, right_current=2):
            last = self.pages
            for num in range(1, last + 1):
                if num <= left_edge or \
                   (num > self.page - left_current - 1 and num < self.page + right_current) or \
                   num > last - right_edge:
                    yield num
                else:
                    yield None
    
    cases_pagination = Pagination(cases, page, per_page, total)
    
    # 获取所有工具（工具不需要分页，因为数量较少）
    tools = sorted(MOCK_TOOLS, key=lambda x: x['rating'], reverse=True)
    
    return render_template('applications.html', cases=cases_pagination, tools=tools)


@main_bp.route('/ethics')
def ethics():
    """伦理与未来模块"""
    record_visit()
    
    # 分页参数
    page = request.args.get('page', 1, type=int)
    per_page = 9
    
    # 获取所有专题
    all_topics = sorted(MOCK_ETHICS_TOPICS, key=lambda x: x['created_at'], reverse=True)
    
    # 分页处理
    total = len(all_topics)
    start = (page - 1) * per_page
    end = start + per_page
    topics = all_topics[start:end]
    
    # 创建分页对象
    class Pagination:
        def __init__(self, items, page, per_page, total):
            self.items = items
            self.page = page
            self.per_page = per_page
            self.total = total
            self.pages = (total + per_page - 1) // per_page if total > 0 else 1
            self.has_prev = page > 1
            self.has_next = page < self.pages
            self.prev_num = page - 1 if self.has_prev else None
            self.next_num = page + 1 if self.has_next else None
        
        def iter_pages(self, left_edge=1, right_edge=1, left_current=2, right_current=2):
            last = self.pages
            for num in range(1, last + 1):
                if num <= left_edge or \
                   (num > self.page - left_current - 1 and num < self.page + right_current) or \
                   num > last - right_edge:
                    yield num
                else:
                    yield None
    
    topics_pagination = Pagination(topics, page, per_page, total)
    
    return render_template('ethics.html', topics=topics_pagination)


@main_bp.route('/ethics/<slug>')
def ethics_topic(slug):
    """伦理专题详情页"""
    topic = next((t for t in MOCK_ETHICS_TOPICS if t['slug'] == slug), None)
    if not topic:
        flash('专题不存在')
        return redirect(url_for('main.ethics'))
    
    # 模拟评论（暂时为空）
    comments = []
    
    return render_template('ethics_topic.html', topic=topic, comments=comments)


@main_bp.route('/resources')
def resources():
    """资源中心模块"""
    record_visit()
    
    # 从数据库获取所有术语（按字母排序）
    all_terms = get_all_terms_from_db()
    terms = sorted(all_terms, key=lambda x: x['term'])
    
    # 分页参数
    page = request.args.get('page', 1, type=int)
    per_page = 9
    
    # 获取推荐阅读资源（排除课程类型）
    all_resources = [r for r in MOCK_RESOURCES if r.get('type') != '课程']
    all_resources = sorted(all_resources, key=lambda x: x['created_at'], reverse=True)
    
    # 获取课程推荐（只包含课程类型）
    all_courses = [r for r in MOCK_RESOURCES if r.get('type') == '课程']
    all_courses = sorted(all_courses, key=lambda x: x['created_at'], reverse=True)
    
    # 分页处理 - 推荐阅读
    total = len(all_resources)
    start = (page - 1) * per_page
    end = start + per_page
    resources = all_resources[start:end]
    
    # 创建分页对象
    class Pagination:
        def __init__(self, items, page, per_page, total):
            self.items = items
            self.page = page
            self.per_page = per_page
            self.total = total
            self.pages = (total + per_page - 1) // per_page if total > 0 else 1
            self.has_prev = page > 1
            self.has_next = page < self.pages
            self.prev_num = page - 1 if self.has_prev else None
            self.next_num = page + 1 if self.has_next else None
        
        def iter_pages(self, left_edge=1, right_edge=1, left_current=2, right_current=2):
            last = self.pages
            for num in range(1, last + 1):
                if num <= left_edge or \
                   (num > self.page - left_current - 1 and num < self.page + right_current) or \
                   num > last - right_edge:
                    yield num
                else:
                    yield None
    
    resources_pagination = Pagination(resources, page, per_page, total)
    
    return render_template('resources.html', 
                          terms=terms, 
                          resources=resources_pagination,
                          courses=all_courses)


@main_bp.route('/community')
def community():
    """社区主页 - 显示社区导航"""
    record_visit()
    return render_template('community.html')


@main_bp.route('/community/forum')
def forum():
    """问答论坛"""
    record_visit()
    
    # 获取搜索关键词
    search_query = request.args.get('q', '').strip()
    page = request.args.get('page', 1, type=int)
    per_page = 10
    
    # 获取所有帖子
    all_posts = MOCK_FORUM_POSTS.copy()
    
    # 如果有搜索关键词，进行过滤
    if search_query:
        search_query_lower = search_query.lower()
        filtered_posts = []
        for post in all_posts:
            # 搜索标题和内容
            title_match = search_query_lower in post.get('title', '').lower()
            content_match = search_query_lower in post.get('content', '').lower()
            if title_match or content_match:
                filtered_posts.append(post)
        all_posts = filtered_posts
    
    # 为帖子添加作者信息和评论数，基于真实数据计算
    posts_list = []
    from app.mock_data import MOCK_COMMENTS, MOCK_LIKES
    for post in all_posts:
        user = get_user_by_id(post['user_id'])
        post_copy = post.copy()
        post_copy['author'] = {'username': user['username']} if user else {'username': '未知用户'}
        # 基于真实评论数据计算评论数
        post_copy['comments_count'] = len(MOCK_COMMENTS.get(post['id'], []))
        # 基于真实点赞数据计算点赞数
        if 'posts' not in MOCK_LIKES:
            MOCK_LIKES['posts'] = {}
        if post['id'] not in MOCK_LIKES['posts']:
            MOCK_LIKES['posts'][post['id']] = []
        post_copy['likes'] = len(MOCK_LIKES['posts'][post['id']])
        # 确保views至少为0
        if 'views' not in post_copy or post_copy['views'] < 0:
            post_copy['views'] = 0
        posts_list.append(post_copy)
    
    # 排序：如果有搜索，按相关性排序；否则按时间排序
    if search_query:
        search_query_lower = search_query.lower()
        def sort_key(post):
            title_match = search_query_lower in post.get('title', '').lower()
            return (not title_match, -post.get('likes', 0), -post.get('views', 0))
        posts_list = sorted(posts_list, key=sort_key)
    else:
        posts_list = sorted(posts_list, key=lambda x: x['created_at'], reverse=True)
    
    # 简单分页
    total = len(posts_list)
    start = (page - 1) * per_page
    end = start + per_page
    posts = posts_list[start:end]
    
    # 创建分页对象（模拟）
    class Pagination:
        def __init__(self, items, page, per_page, total):
            self.items = items
            self.page = page
            self.per_page = per_page
            self.total = total
            self.pages = (total + per_page - 1) // per_page if total > 0 else 1
            self.has_prev = page > 1
            self.has_next = page < self.pages
            self.prev_num = page - 1 if self.has_prev else None
            self.next_num = page + 1 if self.has_next else None
        
        def iter_pages(self, left_edge=1, right_edge=1, left_current=2, right_current=2):
            last = self.pages
            for num in range(1, last + 1):
                if num <= left_edge or \
                   (num > self.page - left_current - 1 and num < self.page + right_current) or \
                   num > last - right_edge:
                    yield num
                else:
                    yield None
    
    pagination = Pagination(posts, page, per_page, total)
    
    # 获取未读消息数
    unread_count = 0
    if current_user.is_authenticated:
        from app.mock_data import MOCK_MESSAGES
        user_messages = MOCK_MESSAGES.get(current_user.id, [])
        unread_count = sum(1 for msg in user_messages if not msg.get('read', False))
    
    return render_template('forum.html', posts=pagination, search_query=search_query, unread_count=unread_count)


@main_bp.route('/community/forum/<int:post_id>')
def forum_post(post_id):
    """论坛帖子详情"""
    from app.mock_data import MOCK_COMMENTS, MOCK_MESSAGES, MOCK_LIKES, get_user_by_id
    
    post = next((p for p in MOCK_FORUM_POSTS if p['id'] == post_id), None)
    if not post:
        flash('帖子不存在')
        return redirect(url_for('main.forum'))
    
    # 增加浏览量
    post['views'] = post.get('views', 0) + 1
    
    user = get_user_by_id(post['user_id'])
    post['author'] = {'username': user['username']} if user else {'username': '未知用户'}
    
    # 基于真实数据更新点赞数
    if 'posts' in MOCK_LIKES and post_id in MOCK_LIKES['posts']:
        post['likes'] = len(MOCK_LIKES['posts'][post_id])
    
    # 获取评论数据，基于真实点赞数据计算点赞数
    comments_data = MOCK_COMMENTS.get(post_id, [])
    comments = []
    for comment_data in comments_data:
        comment_user = get_user_by_id(comment_data['user_id'])
        comment_id = comment_data['id']
        
        # 基于真实点赞数据计算评论点赞数
        if 'comments' not in MOCK_LIKES:
            MOCK_LIKES['comments'] = {}
        if comment_id not in MOCK_LIKES['comments']:
            MOCK_LIKES['comments'][comment_id] = []
        comment_likes = len(MOCK_LIKES['comments'][comment_id])
        
        comment = {
            'id': comment_id,
            'user_id': comment_data['user_id'],
            'author': {'username': comment_user['username']} if comment_user else {'username': '未知用户'},
            'content': comment_data['content'],
            'created_at': comment_data['created_at'],
            'likes': comment_likes,
            'replies': []
        }
        # 处理回复
        for reply_data in comment_data.get('replies', []):
            reply_user = get_user_by_id(reply_data['user_id'])
            reply_id = reply_data['id']
            
            # 基于真实点赞数据计算回复点赞数
            if 'replies' not in MOCK_LIKES:
                MOCK_LIKES['replies'] = {}
            if reply_id not in MOCK_LIKES['replies']:
                MOCK_LIKES['replies'][reply_id] = []
            reply_likes = len(MOCK_LIKES['replies'][reply_id])
            
            comment['replies'].append({
                'id': reply_id,
                'user_id': reply_data['user_id'],
                'author': {'username': reply_user['username']} if reply_user else {'username': '未知用户'},
                'content': reply_data['content'],
                'created_at': reply_data['created_at'],
                'likes': reply_likes
            })
        comments.append(comment)
    
    # 获取未读消息数
    unread_count = 0
    if current_user.is_authenticated:
        user_messages = MOCK_MESSAGES.get(current_user.id, [])
        unread_count = sum(1 for msg in user_messages if not msg.get('read', False))
    
    return render_template('forum_post.html', post=post, comments=comments, unread_count=unread_count)


@main_bp.route('/community/create-post')
@login_required
def create_post():
    """发表帖子页面"""
    return render_template('create_post.html')


@main_bp.route('/community/my-posts')
@login_required
def my_posts():
    """我的帖子"""
    from app.mock_data import MOCK_COMMENTS, MOCK_MESSAGES
    user_posts = [p for p in MOCK_FORUM_POSTS if p['user_id'] == current_user.id]
    # 添加作者信息和评论数
    for post in user_posts:
        user = get_user_by_id(post['user_id'])
        post['author'] = {'username': user['username']} if user else {'username': '未知用户'}
        post['comments_count'] = len(MOCK_COMMENTS.get(post['id'], []))
    
    # 获取未读消息数
    unread_count = 0
    user_messages = MOCK_MESSAGES.get(current_user.id, [])
    unread_count = sum(1 for msg in user_messages if not msg.get('read', False))
    
    return render_template('my_posts.html', posts=user_posts, unread_count=unread_count)


@main_bp.route('/community/my-favorites')
@login_required
def my_favorites():
    """我的收藏"""
    from app.mock_data import USER_FAVORITES, MOCK_COMMENTS, MOCK_MESSAGES
    favorite_ids_dict = USER_FAVORITES.get(current_user.id, {}).get('forum', {})
    favorite_ids = [int(fid) for fid in favorite_ids_dict.keys()]
    favorite_posts = [p for p in MOCK_FORUM_POSTS if p['id'] in favorite_ids]
    # 添加作者信息
    for post in favorite_posts:
        user = get_user_by_id(post['user_id'])
        post['author'] = {'username': user['username']} if user else {'username': '未知用户'}
        post['comments_count'] = len(MOCK_COMMENTS.get(post['id'], []))
    
    # 获取未读消息数
    unread_count = 0
    user_messages = MOCK_MESSAGES.get(current_user.id, [])
    unread_count = sum(1 for msg in user_messages if not msg.get('read', False))
    
    return render_template('my_favorites.html', posts=favorite_posts, unread_count=unread_count)


@main_bp.route('/community/messages')
@login_required
def messages():
    """消息中心"""
    from app.mock_data import MOCK_MESSAGES
    user_messages = MOCK_MESSAGES.get(current_user.id, [])
    # 按时间倒序排列
    user_messages.sort(key=lambda x: x['created_at'], reverse=True)
    # 添加发送者信息
    for msg in user_messages:
        from_user = get_user_by_id(msg['from_user_id'])
        msg['from_user'] = {'username': from_user['username']} if from_user else {'username': '未知用户'}
    unread_count = sum(1 for msg in user_messages if not msg.get('read', False))
    return render_template('messages.html', messages=user_messages, unread_count=unread_count)


@main_bp.route('/community/ai-assistant')
def ai_assistant():
    """AI助教（RAG站内答疑）"""
    return render_template('ai_assistant.html')


@main_bp.route('/admin')
@login_required
def admin_dashboard():
    """管理员后台"""
    if current_user.role not in ['admin', 'super_admin']:
        flash('您没有权限访问此页面')
        return redirect(url_for('main.index'))
    
    # 获取反馈列表（使用全局变量存储）
    from app.routes import feedback_list as feedback_list_data
    
    # 获取访问统计数据
    from app.mock_data import MOCK_VISIT_STATS
    from datetime import date, timedelta
    today = date.today()
    today_visits = MOCK_VISIT_STATS.get(str(today), 0)
    total_visits = sum(MOCK_VISIT_STATS.values())
    
    # 获取最近7天的访问数据
    seven_days_data = []
    for i in range(6, -1, -1):
        visit_date = today - timedelta(days=i)
        date_str = str(visit_date)
        count = MOCK_VISIT_STATS.get(date_str, 0)
        seven_days_data.append({
            'date': visit_date.strftime('%m-%d'),
            'count': count
        })
    
    return render_template('admin_dashboard.html', 
                         feedback_list=feedback_list_data,
                         today_visits=today_visits,
                         total_visits=total_visits,
                         seven_days_data=seven_days_data)


@main_bp.route('/super-admin')
@login_required
def super_admin():
    """超级管理员页面"""
    if current_user.role != 'super_admin':
        flash('您没有权限访问此页面')
        return redirect(url_for('main.index'))
    
    # 从数据库获取所有用户
    try:
        import pymysql
        from app.mock_data import DB_CONFIG
        connection = pymysql.connect(**DB_CONFIG)
        cursor = connection.cursor(pymysql.cursors.DictCursor)
        
        cursor.execute("SELECT * FROM users ORDER BY id")
        db_users = cursor.fetchall()
        cursor.close()
        connection.close()
        
        # 转换为字典格式，添加created_at字段处理
        users = []
        for db_user in db_users:
            user_dict = {
                'id': db_user['id'],
                'username': db_user['username'],
                'email': db_user['email'],
                'role': db_user.get('role', 'user'),
                'created_at': db_user.get('created_at')
            }
            users.append(user_dict)
    except Exception as e:
        # 如果数据库连接失败，使用MOCK_USERS作为后备
        from app.mock_data import MOCK_USERS
        users = sorted(MOCK_USERS, key=lambda x: x['id'])
        print(f"从数据库读取用户失败: {str(e)}，使用MOCK_USERS")
    
    # 获取访问统计数据（与admin_dashboard相同）
    from app.mock_data import MOCK_VISIT_STATS
    from datetime import date, timedelta
    today = date.today()
    today_visits = MOCK_VISIT_STATS.get(str(today), 0)
    total_visits = sum(MOCK_VISIT_STATS.values())
    
    # 获取最近7天的访问数据
    seven_days_data = []
    for i in range(6, -1, -1):
        visit_date = today - timedelta(days=i)
        date_str = str(visit_date)
        count = MOCK_VISIT_STATS.get(date_str, 0)
        seven_days_data.append({
            'date': visit_date.strftime('%m-%d'),
            'count': count
        })
    
    # 获取反馈列表
    from app.routes import feedback_list as feedback_list_data
    
    return render_template('super_admin.html', 
                         users=users,
                         feedback_list=feedback_list_data,
                         today_visits=today_visits,
                         total_visits=total_visits,
                         seven_days_data=seven_days_data)


@main_bp.route('/about')
@main_bp.route('/community/about')
def about():
    """关于我们"""
    return render_template('about.html')


@main_bp.route('/community/privacy')
def privacy():
    """隐私政策"""
    return render_template('privacy.html')


@main_bp.route('/community/guide')
def guide():
    """使用指南"""
    return render_template('guide.html')


@main_bp.route('/search')
def search():
    """站内搜索"""
    from app.category_utils import get_all_categories_for_filter
    
    query = request.args.get('q', '').strip()
    search_type = request.args.get('type', 'general')  # general, advanced, ai
    category = request.args.get('category', '')
    content_type = request.args.get('content_type', '')  # article, tool, term, resource, forum, case, ethics
    time_range = request.args.get('time_range', '')
    
    results = []
    hot_searches = get_hot_searches()  # 获取热门搜索词
    available_categories = get_all_categories_for_filter('article')  # 提供分类选项
    
    if query:
        # 记录搜索日志（用于热门搜索词统计）
        user_id = current_user.id if current_user.is_authenticated else None
        
        if search_type == 'ai':
            # 智能搜索（使用语义理解，使用qwen-max模型）
            results = perform_ai_search(query, category, content_type, time_range)
        elif search_type == 'advanced':
            # 高级检索（带筛选条件）
            results = perform_advanced_search(query, category, content_type, time_range)
        else:
            # 普通搜索
            results = perform_general_search(query, category, content_type, time_range)
        
        # 记录搜索日志
        log_search(query, search_type, user_id, len(results))
    
    return render_template('search.html', 
                         query=query,
                         search_type=search_type,
                         category=category,
                         content_type=content_type,
                         time_range=time_range,
                         results=results,
                         hot_searches=hot_searches,
                         available_categories=available_categories)


def perform_general_search(query, category='', content_type='', time_range=''):
    """执行普通搜索"""
    from app.category_utils import filter_items_by_category
    
    results = []
    query_lower = query.lower()
    now = datetime.now()
    
    # 根据时间范围计算截止日期
    cutoff_date = None
    if time_range == 'week':
        cutoff_date = now - timedelta(days=7)
    elif time_range == 'month':
        cutoff_date = now - timedelta(days=30)
    elif time_range == 'year':
        cutoff_date = now - timedelta(days=365)
    
    # 搜索文章
    if not content_type or content_type == 'article':
        articles_to_search = MOCK_ARTICLES
        if category:
            articles_to_search = filter_items_by_category(articles_to_search, category, 'article')
        
        for article in articles_to_search:
            # 时间筛选
            if cutoff_date and article.get('created_at'):
                if article['created_at'] < cutoff_date:
                    continue
            
            if query_lower in article['title'].lower() or query_lower in article['content'].lower():
                results.append({
                    'type': 'article',
                    'title': article['title'],
                    'content': article['content'][:200] + '...',
                    'url': url_for('main.article_detail', id=article['id']),
                    'date': article.get('created_at', datetime.now())
                })
    
    # 搜索工具
    if not content_type or content_type == 'tool':
        tools_to_search = MOCK_TOOLS
        if category:
            tools_to_search = filter_items_by_category(tools_to_search, category, 'tool')
        
        for tool in tools_to_search:
            # 时间筛选
            if cutoff_date and tool.get('created_at'):
                if tool['created_at'] < cutoff_date:
                    continue
            
            if query_lower in tool['name'].lower() or query_lower in tool.get('description', '').lower():
                results.append({
                    'type': 'tool',
                    'title': tool['name'],
                    'content': tool.get('description', '')[:200] + '...',
                    'url': tool.get('url', '#'),
                    'date': tool.get('created_at', datetime.now())
                })
    
    # 搜索术语
    if not content_type or content_type == 'term':
        terms_to_search = get_all_terms_from_db()
        if category:
            terms_to_search = filter_items_by_category(terms_to_search, category, 'term')
        
        for term in terms_to_search:
            if query_lower in term['term'].lower() or query_lower in term.get('definition', '').lower():
                results.append({
                    'type': 'term',
                    'title': term['term'],
                    'content': term.get('definition', '')[:200] + '...',
                    'url': '#',
                    'date': datetime.now()
                })
    
    # 搜索资源
    if not content_type or content_type == 'resource':
        resources_to_search = MOCK_RESOURCES
        if category:
            resources_to_search = filter_items_by_category(resources_to_search, category, 'resource')
        
        for resource in resources_to_search:
            # 时间筛选
            if cutoff_date and resource.get('created_at'):
                if resource['created_at'] < cutoff_date:
                    continue
            
            if query_lower in resource['title'].lower() or query_lower in resource.get('description', '').lower():
                results.append({
                    'type': 'resource',
                    'title': resource['title'],
                    'content': resource.get('description', '')[:200] + '...',
                    'url': resource.get('url', '#'),
                    'date': resource.get('created_at', datetime.now())
                })
    
    # 搜索论坛帖子
    if not content_type or content_type == 'forum':
        posts_to_search = MOCK_FORUM_POSTS
        if category:
            posts_to_search = [p for p in posts_to_search if p.get('category') == category]
        
        for post in posts_to_search:
            # 时间筛选
            if cutoff_date and post.get('created_at'):
                if post['created_at'] < cutoff_date:
                    continue
            
            if query_lower in post['title'].lower() or query_lower in post.get('content', '').lower():
                results.append({
                    'type': 'forum',
                    'title': post['title'],
                    'content': post.get('content', '')[:200] + '...',
                    'url': url_for('main.forum_post', post_id=post['id']),
                    'date': post.get('created_at', datetime.now())
                })
    
    # 搜索案例
    if not content_type or content_type == 'case':
        cases_to_search = MOCK_CASES
        if category:
            cases_to_search = filter_items_by_category(cases_to_search, category, 'case')
        
        for case in cases_to_search:
            # 时间筛选
            if cutoff_date and case.get('created_at'):
                if case['created_at'] < cutoff_date:
                    continue
            
            if query_lower in case['title'].lower() or query_lower in case.get('description', '').lower():
                results.append({
                    'type': 'case',
                    'title': case['title'],
                    'content': case.get('description', '')[:200] + '...',
                    'url': case.get('external_link', '#'),
                    'date': case.get('created_at', datetime.now())
                })
    
    # 搜索伦理专题
    if not content_type or content_type == 'ethics':
        topics_to_search = MOCK_ETHICS_TOPICS
        if category:
            topics_to_search = filter_items_by_category(topics_to_search, category, 'ethics')
        
        for topic in topics_to_search:
            # 时间筛选
            if cutoff_date and topic.get('created_at'):
                if topic['created_at'] < cutoff_date:
                    continue
            
            if query_lower in topic['title'].lower() or query_lower in topic.get('description', '').lower():
                results.append({
                    'type': 'ethics',
                    'title': topic['title'],
                    'content': topic.get('description', '')[:200] + '...',
                    'url': url_for('main.ethics_topic', slug=topic.get('slug', '')),
                    'date': topic.get('created_at', datetime.now())
                })
    
    return results


def perform_ai_search(query, category='', content_type='', time_range=''):
    """执行智能搜索 - 使用语义理解进行搜索（使用qwen-max模型）"""
    from app.ai_service import ai_search
    
    # 分析搜索意图并提取关键词（使用qwen-max模型）
    keywords = ai_search(query)
    
    if keywords:
        # 使用AI提取的关键词进行搜索，同时保留原始查询
        all_keywords = [query] + keywords
        results = []
        seen = set()
        
        for keyword in all_keywords:
            keyword_results = perform_general_search(keyword, category, content_type, time_range)
            for result in keyword_results:
                result_key = result['title']
                if result_key not in seen:
                    seen.add(result_key)
                    results.append(result)
        
        # 根据相关性排序（标题匹配优先）
        query_lower = query.lower()
        results.sort(key=lambda x: (
            0 if query_lower in x['title'].lower() else 1,  # 标题匹配的优先
            -len(x['title'])  # 然后按标题长度
        ))
        
        return results[:30]  # 限制返回数量
    else:
        # 如果AI搜索失败，回退到普通搜索
        return perform_general_search(query, category, content_type, time_range)


def perform_advanced_search(query, category='', content_type='', time_range=''):
    """执行高级检索 - 带筛选条件的搜索"""
    # 高级检索使用与普通搜索相同的逻辑，但筛选条件更严格
    results = perform_general_search(query, category, content_type, time_range)
    
    # 结果已经根据筛选条件过滤，这里可以添加额外的排序或过滤逻辑
    # 按日期排序（最新的在前）
    results.sort(key=lambda x: x.get('date', datetime.min), reverse=True)
    
    return results


def log_search(query, search_type, user_id, result_count):
    """记录搜索日志"""
    from app.mock_data import SEARCH_LOGS
    if query and len(query.strip()) > 0:
        query_key = query.strip().lower()
        if query_key not in SEARCH_LOGS:
            SEARCH_LOGS[query_key] = {
                'count': 0,
                'last_search': datetime.now(),
                'original_query': query.strip()  # 保留原始查询（保持大小写）
            }
        SEARCH_LOGS[query_key]['count'] += 1
        SEARCH_LOGS[query_key]['last_search'] = datetime.now()
        SEARCH_LOGS[query_key]['result_count'] = result_count


def update_search_log_result_count(query, search_type, result_count):
    """更新搜索日志的结果数量"""
    from app.mock_data import SEARCH_LOGS
    query_key = query.strip().lower()
    if query_key in SEARCH_LOGS:
        SEARCH_LOGS[query_key]['result_count'] = result_count


def get_hot_searches(limit=10):
    """获取热门搜索词"""
    from app.mock_data import SEARCH_LOGS
    
    # 只统计最近30天的搜索
    cutoff_date = datetime.now() - timedelta(days=30)
    
    # 筛选最近30天的搜索，并按搜索次数排序
    recent_searches = [
        {
            'query': log_data['original_query'],
            'count': log_data['count'],
            'last_search': log_data['last_search']
        }
        for query_key, log_data in SEARCH_LOGS.items()
        if log_data.get('last_search', datetime.now()) >= cutoff_date
    ]
    
    # 按搜索次数降序排序
    recent_searches.sort(key=lambda x: x['count'], reverse=True)
    
    return recent_searches[:limit]


@main_bp.route('/article/<int:id>')
def article_detail(id):
    """文章详情页 - 如果有URL则重定向，否则显示简介"""
    article = next((a for a in MOCK_ARTICLES if a['id'] == id), None)
    if not article:
        flash('文章不存在')
        return redirect(url_for('main.index'))
    
    # 如果文章有URL，直接重定向到外部链接
    if article.get('url') and article['url'] != '#':
        return redirect(article['url'], code=302)
    
    # 否则显示简介页面
    article['views'] = article.get('views', 0) + 1
    user = get_user_by_id(article['author_id'])
    if user:
        display_name = user.get('nickname', user['username'])
        article['author'] = {'username': display_name}
    else:
        article['author'] = {'username': '系统'}
    
    # 模拟评论（暂时为空）
    comments = []
    
    return render_template('article_detail.html', article=article, comments=comments)


# 认证蓝图
auth_bp = Blueprint('auth', __name__)

# API蓝图
api_bp = Blueprint('api', __name__)

@auth_bp.route('/login', methods=['GET', 'POST'])
def login():
    """用户登录"""
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username')
        password = data.get('password')
        
        user_data = get_user_by_username(username)
        
        # 使用check_password_hash验证密码
        from werkzeug.security import check_password_hash
        if user_data and check_password_hash(user_data['password'], password):
            user = MockUser(user_data)
            login_user(user, remember=True)
            # 检查是否是首次登录
            first_login = user_data.get('first_login', False)
            # 获取重定向URL
            next_url = session.pop('next_url', None) or url_for('main.index')
            return jsonify({
                'success': True, 
                'message': '登录成功',
                'first_login': first_login,
                'redirect_url': next_url
            })
        else:
            return jsonify({'success': False, 'message': '用户名或密码错误'})
    
    return render_template('login.html')


@auth_bp.route('/register', methods=['GET', 'POST'])
def register():
    """用户注册"""
    from app.mock_data import MOCK_USERS
    import random
    import string
    
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username')
        nickname = data.get('nickname', '').strip()
        email = data.get('email', '').strip()
        password = data.get('password')
        security_question = data.get('security_question', '').strip()
        security_answer = data.get('security_answer', '').strip()
        interests = data.get('interests', [])
        
        # 必填字段验证
        if not username:
            return jsonify({'success': False, 'message': '账号不能为空'})
        if not password:
            return jsonify({'success': False, 'message': '密码不能为空'})
        if not security_question:
            return jsonify({'success': False, 'message': '二级问题不能为空'})
        if not security_answer:
            return jsonify({'success': False, 'message': '问题答案不能为空'})
        
        # 检查用户名是否已存在
        if get_user_by_username(username):
            return jsonify({'success': False, 'message': '账号已存在'})
        
        # 检查邮箱是否已存在（如果提供了邮箱）
        if email:
            for user in MOCK_USERS:
                if user.get('email') == email:
                    return jsonify({'success': False, 'message': '邮箱已被注册'})
        
        # 如果没有提供昵称，随机生成10位英文字符
        if not nickname:
            nickname = ''.join(random.choices(string.ascii_letters, k=10))
        
        # 创建新用户（密码加密）
        from werkzeug.security import generate_password_hash
        new_id = max([u['id'] for u in MOCK_USERS], default=0) + 1
        new_user = {
            'id': new_id,
            'username': username,
            'nickname': nickname,
            'email': email if email else None,
            'password': generate_password_hash(password),
            'role': 'user',
            'avatar': None,
            'security_question': security_question,
            'security_answer': security_answer,
            'interests': interests if interests else [],
            'favorites': {},
            'likes': {},
            'messages': [],
            'first_login': True
        }
        MOCK_USERS.append(new_user)
        
        # 自动登录
        user = MockUser(new_user)
        login_user(user, remember=True)
        
        return jsonify({'success': True, 'message': '注册成功', 'nickname': nickname})
    
    return render_template('register.html')


@api_bp.route('/auth/security-question', methods=['GET'])
def get_security_question_by_username():
    """根据用户名获取二级问题"""
    username = request.args.get('username')
    user_data = get_user_by_username(username)
    if user_data:
        return jsonify({'success': True, 'question': user_data.get('security_question', '')})
    return jsonify({'success': False, 'message': '用户不存在'})


@auth_bp.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    """忘记密码"""
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username')
        security_answer = data.get('security_answer')
        new_password = data.get('new_password')
        
        user_data = get_user_by_username(username)
        if not user_data:
            return jsonify({'success': False, 'message': '用户名不存在'})
        
        # 验证二级问题答案
        if user_data.get('security_answer', '').lower() != security_answer.lower():
            return jsonify({'success': False, 'message': '二级问题答案错误'})
        
        # 更新密码（加密存储）
        from werkzeug.security import generate_password_hash
        user_data['password'] = generate_password_hash(new_password)
        
        return jsonify({'success': True, 'message': '密码重置成功，请重新登录'})
    
    return render_template('forgot_password.html')


@auth_bp.route('/logout')
@login_required
def logout():
    """用户登出"""
    logout_user()
    flash('您已成功登出')
    return redirect(url_for('main.index'))


@api_bp.route('/ai-chat', methods=['POST'])
@login_required
def ai_chat():
    """聊天接口（实验室）- 支持多模型"""
    from app.ai_service import chat_with_model
    
    data = request.get_json()
    message = data.get('message', '')
    model = data.get('model', 'aliyun-qwen-turbo')
    
    if not message:
        return jsonify({
            'success': False,
            'message': '消息不能为空',
            'model': model
        })
    
    # 调用对应模型的API
    response = chat_with_model(message, model)
    return jsonify(response)


@api_bp.route('/forum/post', methods=['POST'])
@login_required
def create_forum_post():
    """创建论坛帖子"""
    from datetime import datetime
    data = request.get_json()
    title = data.get('title', '').strip()
    content = data.get('content', '').strip()
    category = data.get('category', '讨论')
    
    if not title or not content:
        return jsonify({'success': False, 'message': '标题和内容不能为空'})
    
    # 创建新帖子
    new_id = max([p['id'] for p in MOCK_FORUM_POSTS], default=0) + 1
    new_post = {
        'id': new_id,
        'title': title,
        'content': content,
        'user_id': current_user.id,
        'category': category,
        'views': 0,
        'likes': 0,
        'created_at': datetime.now(),
        'comments_count': 0
    }
    MOCK_FORUM_POSTS.append(new_post)
    
    return jsonify({'success': True, 'post_id': new_id, 'message': '帖子发布成功'})


@api_bp.route('/forum/<int:post_id>/comment', methods=['POST'])
@login_required
def add_comment(post_id):
    """添加评论"""
    from app.mock_data import MOCK_COMMENTS, MOCK_MESSAGES, get_user_by_id
    data = request.get_json()
    content = data.get('content', '').strip()
    
    if not content:
        return jsonify({'success': False, 'message': '评论内容不能为空'})
    
    post = next((p for p in MOCK_FORUM_POSTS if p['id'] == post_id), None)
    if not post:
        return jsonify({'success': False, 'message': '帖子不存在'})
    
    # 添加评论到模拟数据
    if post_id not in MOCK_COMMENTS:
        MOCK_COMMENTS[post_id] = []
    
    new_comment_id = max([c['id'] for c in MOCK_COMMENTS[post_id]], default=0) + 1
    new_comment = {
        'id': new_comment_id,
        'user_id': current_user.id,
        'content': content,
        'created_at': datetime.now(),
        'likes': 0,
        'replies': []
    }
    MOCK_COMMENTS[post_id].append(new_comment)
    
    # 更新帖子评论数
    post['comments_count'] = len(MOCK_COMMENTS[post_id])
    
    # 发送消息给帖子作者（如果不是自己）
    if post['user_id'] != current_user.id:
        if post['user_id'] not in MOCK_MESSAGES:
            MOCK_MESSAGES[post['user_id']] = []
        MOCK_MESSAGES[post['user_id']].append({
            'id': len(MOCK_MESSAGES[post['user_id']]) + 1,
            'type': 'reply',
            'content': f'用户"{current_user.username}"评论了您的帖子',
            'related_id': post_id,
            'related_type': 'post',
            'from_user_id': current_user.id,
            'created_at': datetime.now(),
            'read': False
        })
    
    return jsonify({'success': True, 'comment_id': new_comment_id, 'message': '评论发表成功'})


@api_bp.route('/forum/<int:post_id>/comment/<int:comment_id>/like', methods=['POST'])
@login_required
def like_comment(post_id, comment_id):
    """点赞/取消点赞评论"""
    from app.mock_data import MOCK_COMMENTS, MOCK_MESSAGES, MOCK_LIKES
    comment = None
    comments = MOCK_COMMENTS.get(post_id, [])
    for c in comments:
        if c['id'] == comment_id:
            comment = c
            break
    
    if not comment:
        return jsonify({'success': False, 'message': '评论不存在'})
    
    user_id = current_user.id
    if 'comments' not in MOCK_LIKES:
        MOCK_LIKES['comments'] = {}
    if comment_id not in MOCK_LIKES['comments']:
        MOCK_LIKES['comments'][comment_id] = []
    
    is_liked = user_id in MOCK_LIKES['comments'][comment_id]
    
    if is_liked:
        MOCK_LIKES['comments'][comment_id].remove(user_id)
        is_liked = False
    else:
        MOCK_LIKES['comments'][comment_id].append(user_id)
        is_liked = True
        
        # 发送消息给评论作者（如果不是自己）
        if comment.get('user_id') and comment['user_id'] != current_user.id:
            if comment['user_id'] not in MOCK_MESSAGES:
                MOCK_MESSAGES[comment['user_id']] = []
            MOCK_MESSAGES[comment['user_id']].append({
                'id': len(MOCK_MESSAGES[comment['user_id']]) + 1,
                'type': 'like',
                'content': f'用户"{current_user.username}"点赞了您的评论',
                'related_id': comment_id,
                'related_type': 'comment',
                'from_user_id': current_user.id,
                'created_at': datetime.now(),
                'read': False
            })
    
    # 基于真实数据更新点赞数
    comment['likes'] = len(MOCK_LIKES['comments'][comment_id])
    
    return jsonify({
        'success': True,
        'likes': comment['likes'],
        'is_liked': is_liked
    })


@api_bp.route('/forum/<int:post_id>/comment/<int:comment_id>/reply', methods=['POST'])
@login_required
def add_reply(post_id, comment_id):
    """添加回复"""
    from app.mock_data import MOCK_COMMENTS, MOCK_MESSAGES, MOCK_FORUM_POSTS
    data = request.get_json()
    content = data.get('content', '').strip()
    
    if not content:
        return jsonify({'success': False, 'message': '回复内容不能为空'})
    
    comment = None
    comments = MOCK_COMMENTS.get(post_id, [])
    for c in comments:
        if c['id'] == comment_id:
            comment = c
            break
    
    if not comment:
        return jsonify({'success': False, 'message': '评论不存在'})
    
    if 'replies' not in comment:
        comment['replies'] = []
    
    new_reply_id = max([r['id'] for r in comment['replies']], default=0) + 1
    new_reply = {
        'id': new_reply_id,
        'user_id': current_user.id,
        'content': content,
        'created_at': datetime.now(),
        'likes': 0
    }
    comment['replies'].append(new_reply)
    
    # 更新帖子评论数（回复也算作评论）
    post = next((p for p in MOCK_FORUM_POSTS if p['id'] == post_id), None)
    if post:
        post['comments_count'] = len(MOCK_COMMENTS.get(post_id, []))
    
    # 发送消息给评论作者（如果不是自己）
    if comment['user_id'] != current_user.id:
        if comment['user_id'] not in MOCK_MESSAGES:
            MOCK_MESSAGES[comment['user_id']] = []
        MOCK_MESSAGES[comment['user_id']].append({
            'id': len(MOCK_MESSAGES[comment['user_id']]) + 1,
            'type': 'reply',
            'content': f'用户"{current_user.username}"回复了您的评论',
            'related_id': comment_id,
            'related_type': 'comment',
            'from_user_id': current_user.id,
            'created_at': datetime.now(),
            'read': False
        })
    
    return jsonify({'success': True, 'reply_id': new_reply_id, 'message': '回复发表成功'})


@api_bp.route('/forum/<int:post_id>/reply/<int:reply_id>/like', methods=['POST'])
@login_required
def like_reply(post_id, reply_id):
    """点赞/取消点赞回复"""
    from app.mock_data import MOCK_COMMENTS, MOCK_MESSAGES, MOCK_LIKES
    reply = None
    comments = MOCK_COMMENTS.get(post_id, [])
    for c in comments:
        for r in c.get('replies', []):
            if r['id'] == reply_id:
                reply = r
                break
        if reply:
            break
    
    if not reply:
        return jsonify({'success': False, 'message': '回复不存在'})
    
    user_id = current_user.id
    if 'replies' not in MOCK_LIKES:
        MOCK_LIKES['replies'] = {}
    if reply_id not in MOCK_LIKES['replies']:
        MOCK_LIKES['replies'][reply_id] = []
    
    is_liked = user_id in MOCK_LIKES['replies'][reply_id]
    
    if is_liked:
        MOCK_LIKES['replies'][reply_id].remove(user_id)
        is_liked = False
    else:
        MOCK_LIKES['replies'][reply_id].append(user_id)
        is_liked = True
        
        # 发送消息给回复作者（如果不是自己）
        if reply.get('user_id') and reply['user_id'] != current_user.id:
            if reply['user_id'] not in MOCK_MESSAGES:
                MOCK_MESSAGES[reply['user_id']] = []
            MOCK_MESSAGES[reply['user_id']].append({
                'id': len(MOCK_MESSAGES[reply['user_id']]) + 1,
                'type': 'like',
                'content': f'用户"{current_user.username}"点赞了您的回复',
                'related_id': reply_id,
                'related_type': 'reply',
                'from_user_id': current_user.id,
                'created_at': datetime.now(),
                'read': False
            })
    
    # 基于真实数据更新点赞数
    reply['likes'] = len(MOCK_LIKES['replies'][reply_id])
    
    return jsonify({
        'success': True,
        'likes': reply['likes'],
        'is_liked': is_liked
    })


# 注意：已废弃FORUM_LIKES和ETHICS_LIKES，统一使用MOCK_LIKES

@api_bp.route('/forum/<int:post_id>/like', methods=['POST'])
@login_required
def like_post(post_id):
    """点赞/取消点赞帖子"""
    from app.mock_data import MOCK_LIKES, MOCK_MESSAGES
    post = next((p for p in MOCK_FORUM_POSTS if p['id'] == post_id), None)
    if not post:
        return jsonify({'success': False, 'message': '帖子不存在'})
    
    user_id = current_user.id
    
    # 初始化MOCK_LIKES结构
    if 'posts' not in MOCK_LIKES:
        MOCK_LIKES['posts'] = {}
    if post_id not in MOCK_LIKES['posts']:
        MOCK_LIKES['posts'][post_id] = []
    
    # 检查当前点赞状态（基于MOCK_LIKES）
    is_liked = user_id in MOCK_LIKES['posts'][post_id]
    
    # 切换点赞状态
    if is_liked:
        # 取消点赞
        if user_id in MOCK_LIKES['posts'][post_id]:
            MOCK_LIKES['posts'][post_id].remove(user_id)
        is_liked = False
    else:
        # 点赞
        if user_id not in MOCK_LIKES['posts'][post_id]:
            MOCK_LIKES['posts'][post_id].append(user_id)
        is_liked = True
        
        # 发送消息给帖子作者（如果不是自己）
        if post['user_id'] != current_user.id:
            if post['user_id'] not in MOCK_MESSAGES:
                MOCK_MESSAGES[post['user_id']] = []
            MOCK_MESSAGES[post['user_id']].append({
                'id': len(MOCK_MESSAGES[post['user_id']]) + 1,
                'type': 'like',
                'content': f'用户"{current_user.username}"点赞了您的帖子',
                'related_id': post_id,
                'related_type': 'post',
                'from_user_id': current_user.id,
                'created_at': datetime.now(),
                'read': False
            })
    
    # 更新帖子点赞数（基于真实数据）
    post['likes'] = len(MOCK_LIKES['posts'][post_id])
    
    return jsonify({
        'success': True, 
        'likes': post['likes'],
        'is_liked': is_liked
    })


@api_bp.route('/ethics/<int:topic_id>/like', methods=['POST'])
@login_required
def like_ethics_topic(topic_id):
    """点赞/取消点赞伦理专题"""
    from app.mock_data import USER_LIKES
    topic = next((t for t in MOCK_ETHICS_TOPICS if t['id'] == topic_id), None)
    if not topic:
        return jsonify({'success': False, 'message': '专题不存在'})
    
    user_id = current_user.id
    
    # 初始化用户点赞记录（使用USER_LIKES）
    if user_id not in USER_LIKES:
        USER_LIKES[user_id] = {'ethics': {}}
    if 'ethics' not in USER_LIKES[user_id]:
        USER_LIKES[user_id]['ethics'] = {}
    
    # 检查当前点赞状态
    is_liked = str(topic_id) in USER_LIKES[user_id]['ethics'] and USER_LIKES[user_id]['ethics'][str(topic_id)] is not None
    
    # 切换点赞状态
    if is_liked:
        # 取消点赞
        topic['likes'] = max(0, topic.get('likes', 0) - 1)
        USER_LIKES[user_id]['ethics'][str(topic_id)] = None  # 删除记录
        is_liked = False
    else:
        # 点赞
        topic['likes'] = topic.get('likes', 0) + 1
        USER_LIKES[user_id]['ethics'][str(topic_id)] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        is_liked = True
    
    return jsonify({
        'success': True, 
        'likes': topic['likes'],
        'is_liked': is_liked
    })


# 收藏功能API
@api_bp.route('/article/<int:article_id>/favorite', methods=['POST'])
@login_required
def favorite_article(article_id):
    """收藏/取消收藏文章"""
    from app.mock_data import USER_FAVORITES
    article = next((a for a in MOCK_ARTICLES if a['id'] == article_id), None)
    if not article:
        return jsonify({'success': False, 'message': '文章不存在'})
    
    user_id = current_user.id
    if user_id not in USER_FAVORITES:
        USER_FAVORITES[user_id] = {}
    if 'article' not in USER_FAVORITES[user_id]:
        USER_FAVORITES[user_id]['article'] = {}
    
    # 切换收藏状态
    is_favorited = article_id in USER_FAVORITES[user_id]['article']
    if is_favorited:
        del USER_FAVORITES[user_id]['article'][article_id]
        is_favorited = False
    else:
        USER_FAVORITES[user_id]['article'][article_id] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        is_favorited = True
    
    return jsonify({'success': True, 'is_favorited': is_favorited})


@api_bp.route('/forum/<int:post_id>/favorite', methods=['POST'])
@login_required
def favorite_forum_post(post_id):
    """收藏/取消收藏论坛帖子"""
    from app.mock_data import USER_FAVORITES
    post = next((p for p in MOCK_FORUM_POSTS if p['id'] == post_id), None)
    if not post:
        return jsonify({'success': False, 'message': '帖子不存在'})
    
    user_id = current_user.id
    if user_id not in USER_FAVORITES:
        USER_FAVORITES[user_id] = {}
    if 'forum' not in USER_FAVORITES[user_id]:
        USER_FAVORITES[user_id]['forum'] = {}
    
    # 切换收藏状态
    is_favorited = post_id in USER_FAVORITES[user_id]['forum']
    if is_favorited:
        del USER_FAVORITES[user_id]['forum'][post_id]
        is_favorited = False
    else:
        USER_FAVORITES[user_id]['forum'][post_id] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        is_favorited = True
    
    return jsonify({'success': True, 'is_favorited': is_favorited})


@api_bp.route('/forum/<int:post_id>/status', methods=['GET'])
@login_required
def get_post_status(post_id):
    """获取帖子的点赞和收藏状态"""
    from app.mock_data import USER_FAVORITES, MOCK_LIKES
    user_id = current_user.id
    
    # 检查点赞状态
    is_liked = False
    if 'posts' in MOCK_LIKES and post_id in MOCK_LIKES['posts']:
        is_liked = user_id in MOCK_LIKES['posts'][post_id]
    
    # 检查收藏状态
    is_favorited = False
    if user_id in USER_FAVORITES and 'forum' in USER_FAVORITES[user_id]:
        is_favorited = post_id in USER_FAVORITES[user_id]['forum']
    
    return jsonify({
        'success': True,
        'is_liked': is_liked,
        'is_favorited': is_favorited
    })


@api_bp.route('/ethics/<int:topic_id>/favorite', methods=['POST'])
@login_required
def favorite_ethics_topic(topic_id):
    """收藏/取消收藏伦理专题"""
    from app.mock_data import USER_FAVORITES
    topic = next((t for t in MOCK_ETHICS_TOPICS if t['id'] == topic_id), None)
    if not topic:
        return jsonify({'success': False, 'message': '专题不存在'})
    
    user_id = current_user.id
    if user_id not in USER_FAVORITES:
        USER_FAVORITES[user_id] = {}
    if 'ethics' not in USER_FAVORITES[user_id]:
        USER_FAVORITES[user_id]['ethics'] = {}
    
    # 切换收藏状态
    is_favorited = topic_id in USER_FAVORITES[user_id]['ethics']
    if is_favorited:
        del USER_FAVORITES[user_id]['ethics'][topic_id]
        is_favorited = False
    else:
        USER_FAVORITES[user_id]['ethics'][topic_id] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        is_favorited = True
    
    return jsonify({'success': True, 'is_favorited': is_favorited})


@api_bp.route('/ethics/<int:topic_id>/comment', methods=['POST'])
@login_required
def add_ethics_comment(topic_id):
    """添加伦理专题评论"""
    data = request.get_json()
    # 模拟添加评论
    return jsonify({'success': True, 'comment_id': 1, 'message': '评论发表成功（模拟数据）'})


@api_bp.route('/models', methods=['GET'])
def get_models():
    """获取可用模型列表"""
    from app.ai_service import get_available_models
    models = get_available_models()
    return jsonify({'success': True, 'models': models})


# 个人中心路由
@main_bp.route('/profile')
@login_required
def profile():
    """个人中心"""
    return render_template('profile.html')


# 个人中心API
@api_bp.route('/profile/favorites', methods=['GET'])
@login_required
def get_favorites():
    """获取收藏列表"""
    from app.mock_data import USER_FAVORITES, MOCK_ARTICLES, MOCK_FORUM_POSTS, MOCK_ETHICS_TOPICS
    user_id = current_user.id
    favorites = USER_FAVORITES.get(user_id, {})
    
    result = []
    # 处理文章收藏
    if 'article' in favorites and isinstance(favorites['article'], dict):
        for article_id, timestamp in favorites['article'].items():
            if timestamp:  # 只显示已收藏的
                article = next((a for a in MOCK_ARTICLES if a['id'] == int(article_id)), None)
                if article:
                    result.append({
                        'type': '文章',
                        'title': article['title'],
                        'url': url_for('main.article_detail', id=article['id']),
                        'timestamp': timestamp
                    })
    
    # 处理论坛帖子收藏
    if 'forum' in favorites and isinstance(favorites['forum'], dict):
        for post_id, timestamp in favorites['forum'].items():
            if timestamp:  # 只显示已收藏的
                post = next((p for p in MOCK_FORUM_POSTS if p['id'] == int(post_id)), None)
                if post:
                    result.append({
                        'type': '论坛帖子',
                        'title': post['title'],
                        'url': url_for('main.forum_post', post_id=post['id']),
                        'timestamp': timestamp
                    })
    
    # 处理伦理专题收藏
    if 'ethics' in favorites and isinstance(favorites['ethics'], dict):
        for topic_id, timestamp in favorites['ethics'].items():
            if timestamp:  # 只显示已收藏的
                topic = next((t for t in MOCK_ETHICS_TOPICS if t['id'] == int(topic_id)), None)
                if topic:
                    result.append({
                        'type': '伦理专题',
                        'title': topic['title'],
                        'url': url_for('main.ethics_topic', slug=topic['slug']),
                        'timestamp': timestamp
                    })
    
    # 按时间倒序排列
    result.sort(key=lambda x: x.get('timestamp', ''), reverse=True)
    
    return jsonify({'success': True, 'favorites': result})


@api_bp.route('/profile/likes', methods=['GET'])
@login_required
def get_likes():
    """获取点赞记录"""
    from app.mock_data import USER_LIKES, MOCK_FORUM_POSTS, MOCK_ETHICS_TOPICS
    
    user_id = current_user.id
    likes = USER_LIKES.get(user_id, {})
    
    result = []
    # 论坛帖子点赞
    if 'forum' in likes and isinstance(likes['forum'], dict):
        for post_id, timestamp in likes['forum'].items():
            if timestamp:  # 只显示已点赞的
                post = next((p for p in MOCK_FORUM_POSTS if p['id'] == int(post_id)), None)
                if post:
                    result.append({
                        'type': '论坛帖子',
                        'title': post['title'],
                        'url': url_for('main.forum_post', post_id=post['id']),
                        'timestamp': timestamp
                    })
    
    # 伦理专题点赞
    if 'ethics' in likes and isinstance(likes['ethics'], dict):
        for topic_id, timestamp in likes['ethics'].items():
            if timestamp:  # 只显示已点赞的
                topic = next((t for t in MOCK_ETHICS_TOPICS if t['id'] == int(topic_id)), None)
                if topic:
                    result.append({
                        'type': '伦理专题',
                        'title': topic['title'],
                        'url': url_for('main.ethics_topic', slug=topic['slug']),
                        'timestamp': timestamp
                    })
    
    # 按时间倒序排列
    result.sort(key=lambda x: x.get('timestamp', ''), reverse=True)
    
    return jsonify({'success': True, 'likes': result})


@api_bp.route('/profile/messages', methods=['GET'])
@login_required
def get_messages():
    """获取消息列表"""
    from app.mock_data import USER_MESSAGES
    user_id = current_user.id
    messages = USER_MESSAGES.get(user_id, [])
    # 按时间倒序排列
    messages.sort(key=lambda x: x.get('timestamp', ''), reverse=True)
    return jsonify({'success': True, 'messages': messages})


@api_bp.route('/profile/messages/<int:message_id>/read', methods=['POST'])
@login_required
def mark_message_read(message_id):
    """标记消息为已读"""
    from app.mock_data import MOCK_MESSAGES
    user_id = current_user.id
    messages = MOCK_MESSAGES.get(user_id, [])
    for msg in messages:
        if msg.get('id') == message_id:
            msg['read'] = True
            break
    return jsonify({'success': True})


@api_bp.route('/auth/check-first-login', methods=['GET'])
@login_required
def check_first_login():
    """检查是否首次登录"""
    from app.mock_data import get_user_by_id
    user_data = get_user_by_id(current_user.id)
    if user_data:
        return jsonify({'success': True, 'first_login': user_data.get('first_login', False)})
    return jsonify({'success': False})


@api_bp.route('/profile/interests', methods=['POST'])
@login_required
def update_interests():
    """更新用户兴趣"""
    from app.mock_data import get_user_by_id
    data = request.get_json()
    interests = data.get('interests', [])
    
    user_data = get_user_by_id(current_user.id)
    if user_data:
        user_data['interests'] = interests
        user_data['first_login'] = False  # 标记已设置兴趣
        return jsonify({'success': True})
    return jsonify({'success': False, 'message': '用户不存在'})


@api_bp.route('/profile/security-question', methods=['GET'])
@login_required
def get_security_question():
    """获取二级问题"""
    from app.mock_data import get_user_by_id
    user_data = get_user_by_id(current_user.id)
    if user_data:
        return jsonify({'success': True, 'question': user_data.get('security_question', '')})
    return jsonify({'success': False, 'message': '用户不存在'})


@api_bp.route('/profile/password', methods=['POST'])
@login_required
def change_password():
    """修改密码"""
    from app.mock_data import get_user_by_id, MOCK_USERS
    data = request.get_json()
    security_answer = data.get('security_answer', '')
    new_password = data.get('new_password', '')
    
    user_data = get_user_by_id(current_user.id)
    if not user_data:
        return jsonify({'success': False, 'message': '用户不存在'})
    
    # 验证二级问题答案
    if user_data.get('security_answer', '').lower() != security_answer.lower():
        return jsonify({'success': False, 'message': '二级问题答案错误'})
    
        # 更新密码（加密存储）
        from werkzeug.security import generate_password_hash
        user_data['password'] = generate_password_hash(new_password)
        # 强制退出（不提醒用户）
    from flask_login import logout_user
    logout_user()
    
    return jsonify({'success': True, 'message': '密码修改成功，请重新登录'})


@api_bp.route('/profile/security', methods=['POST'])
@login_required
def change_security():
    """修改二级密码"""
    from app.mock_data import get_user_by_id
    data = request.get_json()
    current_answer = data.get('current_answer', '')
    new_question = data.get('new_question', '')
    new_answer = data.get('new_answer', '')
    
    user_data = get_user_by_id(current_user.id)
    if not user_data:
        return jsonify({'success': False, 'message': '用户不存在'})
    
    # 验证当前二级问题答案
    if user_data.get('security_answer', '').lower() != current_answer.lower():
        return jsonify({'success': False, 'message': '当前二级问题答案错误'})
    
    # 更新二级问题
    user_data['security_question'] = new_question
    user_data['security_answer'] = new_answer
    
    return jsonify({'success': True, 'message': '二级密码修改成功'})


@api_bp.route('/profile/username', methods=['POST'])
@login_required
def update_username():
    """更新昵称"""
    from app.mock_data import get_user_by_id, get_user_by_username
    data = request.get_json()
    new_username = data.get('username', '').strip()
    
    if not new_username:
        return jsonify({'success': False, 'message': '昵称不能为空'})
    
    # 检查用户名是否已存在
    existing_user = get_user_by_username(new_username)
    if existing_user and existing_user['id'] != current_user.id:
        return jsonify({'success': False, 'message': '用户名已存在'})
    
    user_data = get_user_by_id(current_user.id)
    if user_data:
        user_data['username'] = new_username
        return jsonify({'success': True, 'message': '昵称更新成功'})
    
    return jsonify({'success': False, 'message': '用户不存在'})


@api_bp.route('/profile/avatar', methods=['POST'])
@login_required
def upload_avatar():
    """上传头像"""
    from werkzeug.utils import secure_filename
    import os
    from app.mock_data import get_user_by_id
    
    if 'avatar' not in request.files:
        return jsonify({'success': False, 'message': '没有上传文件'})
    
    file = request.files['avatar']
    if file.filename == '':
        return jsonify({'success': False, 'message': '没有选择文件'})
    
    # 检查文件类型
    if not file.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
        return jsonify({'success': False, 'message': '只支持图片格式'})
    
    # 保存文件
    upload_folder = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'static', 'uploads', 'avatars')
    os.makedirs(upload_folder, exist_ok=True)
    
    filename = secure_filename(f"{current_user.id}_{datetime.now().strftime('%Y%m%d%H%M%S')}_{file.filename}")
    filepath = os.path.join(upload_folder, filename)
    file.save(filepath)
    
    # 更新用户头像
    user_data = get_user_by_id(current_user.id)
    if user_data:
        # 使用相对路径，前端会通过url_for处理
        avatar_path = f'uploads/avatars/{filename}'
        user_data['avatar'] = avatar_path
        # 返回完整URL路径
        from flask import url_for
        avatar_url = url_for('static', filename=avatar_path)
        return jsonify({'success': True, 'avatar_url': avatar_url})
    
    return jsonify({'success': False, 'message': '用户不存在'})


@api_bp.route('/ai-assistant', methods=['POST'])
@login_required
def ai_assistant_api():
    """智能助教接口 - 基于系统资源的问答"""
    from app.ai_service import chat_with_model
    from config import Config
    from app.mock_data import (
        MOCK_ARTICLES, MOCK_TERMS, MOCK_TOOLS, MOCK_CASES, 
        MOCK_RESOURCES, MOCK_ETHICS_TOPICS, MOCK_LEARNING_PATHS
    )
    
    data = request.get_json()
    question = data.get('question', '')
    
    if not question:
        return jsonify({
            'success': False,
            'message': '问题不能为空'
        })
    
    # RAG检索：从系统资源中查找相关内容（扩展到整个系统）
    question_lower = question.lower()
    relevant_content = []
    sources = []
    
    try:
        from app import db
        # 1. 检索数据库中的文章
        from app.models import Article
        articles = Article.query.filter(
            db.or_(
                Article.title.ilike(f'%{question}%'),
                Article.content.ilike(f'%{question}%')
            )
        ).limit(5).all()
        for article in articles:
            content_preview = article.content[:300] if article.content else ''
            relevant_content.append(f"文章《{article.title}》：{content_preview}...")
            sources.append(f"文章：{article.title}")
        
        # 2. 检索数据库中的术语
        all_terms = get_all_terms_from_db()
        for term in all_terms:
            if (question_lower in term['term'].lower() or 
                question_lower in term['definition'].lower()):
                relevant_content.append(f"术语【{term['term']}】：{term['definition']}")
                sources.append(f"AI术语：{term['term']}")
        
        # 3. 检索数据库中的论坛帖子
        from app.models import ForumPost
        forum_posts = ForumPost.query.filter(
            db.or_(
                ForumPost.title.ilike(f'%{question}%'),
                ForumPost.content.ilike(f'%{question}%')
            )
        ).limit(5).all()
        for post in forum_posts:
            content_preview = post.content[:300] if post.content else ''
            relevant_content.append(f"论坛帖子《{post.title}》：{content_preview}...")
            sources.append(f"问答论坛：{post.title}")
        
        # 4. 检索数据库中的评论
        from app.models import Comment
        comments = Comment.query.filter(
            Comment.content.ilike(f'%{question}%')
        ).limit(5).all()
        for comment in comments:
            content_preview = comment.content[:200] if comment.content else ''
            relevant_content.append(f"评论：{content_preview}...")
            sources.append(f"用户评论")
        
        # 5. 检索数据库中的资源
        from app.models import Resource
        resources = Resource.query.filter(
            db.or_(
                Resource.title.ilike(f'%{question}%'),
                Resource.description.ilike(f'%{question}%')
            )
        ).limit(5).all()
        for resource in resources:
            desc = resource.description[:200] if resource.description else ''
            relevant_content.append(f"资源《{resource.title}》：{desc}")
            sources.append(f"推荐阅读：{resource.title}")
        
        # 6. 检索数据库中的伦理专题
        from app.models import EthicsTopic
        ethics_topics = EthicsTopic.query.filter(
            db.or_(
                EthicsTopic.title.ilike(f'%{question}%'),
                EthicsTopic.description.ilike(f'%{question}%')
            )
        ).limit(5).all()
        for topic in ethics_topics:
            desc = topic.description[:200] if topic.description else ''
            relevant_content.append(f"专题《{topic.title}》：{desc}")
            sources.append(f"伦理与未来：{topic.title}")
        
        # 7. 检索数据库中的工具
        from app.models import Tool
        tools = Tool.query.filter(
            db.or_(
                Tool.name.ilike(f'%{question}%'),
                Tool.description.ilike(f'%{question}%')
            )
        ).limit(5).all()
        for tool in tools:
            desc = tool.description[:200] if tool.description else ''
            relevant_content.append(f"工具《{tool.name}》：{desc}")
            sources.append(f"AI工具：{tool.name}")
        
        # 8. 检索数据库中的案例
        from app.models import Case
        cases = Case.query.filter(
            db.or_(
                Case.title.ilike(f'%{question}%'),
                Case.description.ilike(f'%{question}%')
            )
        ).limit(5).all()
        for case in cases:
            desc = case.description[:200] if case.description else ''
            relevant_content.append(f"案例《{case.title}》：{desc}")
            sources.append(f"应用案例：{case.title}")
        
        # 9. 检索Mock数据作为补充（如果数据库中没有足够内容）
        if len(relevant_content) < 3:
            # 检索Mock文章
            for article in MOCK_ARTICLES:
                if (question_lower in article['title'].lower() or 
                    question_lower in article['content'].lower()[:500]):
                    relevant_content.append(f"文章《{article['title']}》：{article['content'][:300]}...")
                    sources.append(f"文章：{article['title']}")
            
            # 检索Mock工具
            for tool in MOCK_TOOLS:
                if (question_lower in tool['name'].lower() or 
                    question_lower in tool['description'].lower()):
                    relevant_content.append(f"工具《{tool['name']}》：{tool['description']}")
                    sources.append(f"AI工具：{tool['name']}")
            
            # 检索Mock案例
            for case in MOCK_CASES:
                if (question_lower in case['title'].lower() or 
                    question_lower in case['description'].lower()):
                    relevant_content.append(f"案例《{case['title']}》：{case['description']}")
                    sources.append(f"应用案例：{case['title']}")
            
            # 检索Mock资源
            for resource in MOCK_RESOURCES:
                if (question_lower in resource['title'].lower() or 
                    question_lower in resource.get('description', '').lower()):
                    relevant_content.append(f"资源《{resource['title']}》：{resource.get('description', '')[:200]}")
                    sources.append(f"推荐阅读：{resource['title']}")
            
            # 检索Mock伦理专题
            for topic in MOCK_ETHICS_TOPICS:
                if (question_lower in topic['title'].lower() or 
                    question_lower in topic.get('description', '').lower()):
                    relevant_content.append(f"专题《{topic['title']}》：{topic.get('description', '')[:200]}")
                    sources.append(f"伦理与未来：{topic['title']}")
    except Exception as e:
        print(f'RAG检索错误: {str(e)}')
        # 如果数据库查询失败，回退到Mock数据
        for article in MOCK_ARTICLES:
            if (question_lower in article['title'].lower() or 
                question_lower in article['content'].lower()[:500]):
                relevant_content.append(f"文章《{article['title']}》：{article['content'][:300]}...")
                sources.append(f"文章：{article['title']}")
    
    # 构建RAG提示词
    if relevant_content:
        context = "\n\n".join(relevant_content[:10])  # 最多使用10个相关片段（扩展到整个系统后内容更多）
        rag_prompt = f"""你是一个专业的学习助手，专门回答相关问题。

用户问题：{question}

以下是系统内检索到的相关内容：
{context}

请基于以上系统内的实际信息来回答用户的问题。如果系统内的信息不足以完全回答问题，可以结合你的知识进行补充，但请明确标注哪些是系统内的信息，哪些是补充说明。

要求：
1. 优先使用系统内检索到的信息
2. 回答要专业但易懂，适合不同背景的学习者
3. 如果系统内没有相关信息，请直接回复"抱歉，暂时未查找到相关内容"
4. 回答要准确、详细，并提供清晰的解释和示例"""
    else:
        # 没有找到相关内容
        rag_prompt = f"""你是一个专业的学习助手，专门回答相关问题。

用户问题：{question}

系统内未检索到相关内容。

请直接回复："抱歉，暂时未查找到相关内容。"
不要使用你的知识库来回答，因为系统要求只能基于系统内的实际信息来回答。"""
    
    # 优先使用qwen-max，如果没有配置则尝试其他模型
    model = 'aliyun-qwen-max'
    if not Config.DASHSCOPE_API_KEY:
        # 按优先级尝试其他模型
        if Config.DEEPSEEK_API_KEY:
            model = 'deepseek-chat'
        elif Config.KIMI_API_KEY:
            model = 'kimi-moonshot-v1-8k'
        elif Config.GEMINI_API_KEY:
            model = 'gemini-1.5-flash'
        elif Config.OPENAI_API_KEY:
            model = 'openai-gpt-3.5-turbo'
    
    # 调用API
    response = chat_with_model(rag_prompt, model)
    
    if response['success']:
        answer = response['message']
        # 检查是否是没有找到内容
        no_content_found = '抱歉，暂时未查找到相关内容' in answer or '未检索到' in answer
        
        return jsonify({
            'success': True,
            'answer': answer,
            'sources': sources[:3] if sources else [],  # 最多返回3个来源
            'no_content_found': no_content_found
        })
    else:
        return jsonify({
            'success': False,
            'message': response['message']
        })


@api_bp.route('/feedback', methods=['POST'])
@login_required
def submit_feedback():
    """提交用户反馈"""
    data = request.get_json()
    feedback_type = data.get('type', 'general')  # ai_assistant, general等
    question = data.get('question', '')
    content = data.get('content', '').strip()
    
    if not content:
        return jsonify({
            'success': False,
            'message': '反馈内容不能为空'
        })
    
    # 存储反馈（使用全局变量，实际应保存到数据库）
    from app.routes import feedback_list
    feedback_item = {
        'id': len(feedback_list) + 1,
        'user_id': current_user.id,
        'username': current_user.username,
        'type': feedback_type,
        'question': question,
        'content': content,
        'created_at': datetime.now(),
        'status': 'pending'  # pending, resolved
    }
    feedback_list.append(feedback_item)
    
    return jsonify({
        'success': True,
        'message': '反馈提交成功'
    })


@api_bp.route('/super-admin/users/<int:user_id>/role', methods=['POST'])
@login_required
def change_user_role(user_id):
    """修改用户身份（仅超级管理员）"""
    if current_user.role != 'super_admin':
        return jsonify({
            'success': False,
            'message': '您没有权限执行此操作'
        })
    
    data = request.get_json()
    new_role = data.get('role', '').strip()
    
    if new_role not in ['user', 'admin']:
        return jsonify({
            'success': False,
            'message': '无效的身份类型'
        })
    
    # 修改用户身份
    from app.mock_data import MOCK_USERS
    user = next((u for u in MOCK_USERS if u['id'] == user_id), None)
    if not user:
        return jsonify({
            'success': False,
            'message': '用户不存在'
        })
    
    if user['role'] == 'super_admin':
        return jsonify({
            'success': False,
            'message': '无法修改超级管理员的身份'
        })
    
    user['role'] = new_role
    
    return jsonify({
        'success': True,
        'message': f'用户身份已修改为{new_role}'
    })


@api_bp.route('/admin/feedback/<int:feedback_id>/resolve', methods=['POST'])
@login_required
def resolve_feedback(feedback_id):
    """标记反馈为已处理"""
    if current_user.role not in ['admin', 'super_admin']:
        return jsonify({
            'success': False,
            'message': '您没有权限执行此操作'
        })
    
    # 获取反馈列表
    from app.routes import feedback_list as feedback_list_data
    for feedback in feedback_list_data:
        if feedback['id'] == feedback_id:
            feedback['status'] = 'resolved'
            return jsonify({
                'success': True,
                'message': '反馈已标记为已处理'
            })
    
    return jsonify({
        'success': False,
        'message': '反馈不存在'
    })


@api_bp.route('/playground/<tool_type>', methods=['POST'])
@login_required
def playground_tool(tool_type):
    """工具集接口 - 每个工具使用专门的助手"""
    from app.ai_service import chat_with_model
    from config import Config
    
    data = request.get_json()
    user_input = data.get('input', '')
    model = data.get('model', 'aliyun-qwen-turbo')
    conversation = data.get('conversation', [])  # 对话历史（仅编程工具使用）
    
    if not user_input:
        return jsonify({
            'success': False,
            'message': '输入内容不能为空'
        })
    
    # 为每个工具定义专门的agent提示词
    agent_prompts = {
        'image-gen': """你是一个专业的图像生成助手。用户会描述他们想要生成的图像，你需要：
1. 理解用户的描述
2. 优化和扩展描述，使其更适合图像生成
3. 提供详细的图像生成提示词（prompt）

注意：你只负责生成提示词，不直接生成图像。请用中文回复。""",
        'writing': """你是一个专业的写作助手。你的任务是帮助用户进行各种写作工作，包括：
- 文章写作
- 文案创作
- 内容优化
- 创意写作

请根据用户的需求，提供高质量的写作内容。用中文回复。""",
        'translation': """你是一个专业的翻译助手。你的任务是：
1. 准确翻译用户提供的文本
2. 保持原文的语气和风格
3. 提供多种翻译选项（如果适用）

请用中文回复。""",
        'programming': """你是一个专业的编程助手。你的任务是：
1. 帮助用户编写代码
2. 调试和优化代码
3. 解释代码逻辑
4. 提供编程建议

请用中文回复，代码部分使用代码块格式（```语言名\n代码\n```）。如果用户的问题涉及代码，请提供完整的、可运行的代码。""",
        'ppt': """你是一个专业的PPT制作助手。你的任务是：
1. 根据用户的需求，生成PPT的大纲和内容
2. 为每个幻灯片提供标题和要点
3. 提供结构化的内容，方便制作PPT
4. 使用Markdown格式输出，便于后续处理

输出格式：
# 幻灯片1标题
- 要点1
- 要点2

# 幻灯片2标题
- 要点1
- 要点2

请用中文回复。"""
    }
    
    # 特殊处理：图片生成直接调用VolcEngine Seedream API
    if tool_type == 'image-gen':
        return generate_image_with_seedream(user_input)
    
    # 获取对应工具的agent提示词
    system_prompt = agent_prompts.get(tool_type, '你是一个专业的助手。请用中文回复。')
    
    # 如果是编程工具且有对话历史，构建多轮对话
    if tool_type == 'programming' and conversation:
        # 构建对话消息列表
        messages = [{'role': 'system', 'content': system_prompt}]
        # 添加历史对话
        for msg in conversation:
            messages.append({'role': msg['role'], 'content': msg['content']})
        # 添加当前用户输入
        messages.append({'role': 'user', 'content': user_input})
        
        # 调用支持多轮对话的API
        response = chat_with_model_multi_turn(messages, model)
    else:
        # 单轮对话
        full_message = f"{system_prompt}\n\n用户需求：{user_input}"
        response = chat_with_model(full_message, model)
    
    if response['success']:
        result = {
            'success': True,
            'output': response['message'],
            'model': response.get('model', model)
        }
        
        # 如果是PPT制作，标记为PPT格式
        if tool_type == 'ppt':
            result['type'] = 'ppt'
        
        # 如果是编程工具，识别代码语言
        if tool_type == 'programming':
            language = detect_code_language(response['message'])
            result['language'] = language
        
        # 保存AI使用历史记录（除了图片生成，图片生成在generate_image_with_seedream中处理）
        from flask_login import current_user
        if current_user.is_authenticated and tool_type != 'image-gen':
            save_ai_usage_history(
                user_id=current_user.id,
                tool_type=tool_type,
                input_text=user_input[:2000] if user_input else None,
                output_text=response['message'][:5000] if response['message'] else None,
                model_used=response.get('model', model)
            )
            
            # 保存对话历史（支持多轮对话）
            conversation_id = data.get('conversation_id')  # 如果有conversation_id，更新现有记录
            if tool_type == 'programming' and conversation:
                # 添加当前对话到历史
                updated_conversation = conversation + [
                    {'role': 'user', 'content': user_input},
                    {'role': 'assistant', 'content': response['message']}
                ]
                save_conversation_history(
                    user_id=current_user.id,
                    tool_type=tool_type,
                    conversation_data=updated_conversation,
                    model_used=response.get('model', model),
                    conversation_id=conversation_id
                )
                # 返回conversation_id以便前端继续使用
                result['conversation_id'] = get_or_create_conversation_id(
                    current_user.id, tool_type, conversation_id
                )
            elif tool_type != 'programming':
                # 非编程工具也保存单次对话历史
                conversation_data = [
                    {'role': 'user', 'content': user_input},
                    {'role': 'assistant', 'content': response['message']}
                ]
                save_conversation_history(
                    user_id=current_user.id,
                    tool_type=tool_type,
                    conversation_data=conversation_data,
                    model_used=response.get('model', model)
                )
        
        return jsonify(result)
    else:
        return jsonify({
            'success': False,
            'message': response['message']
        })


def generate_image_with_seedream(prompt):
    """使用火山引擎Seedream API生成图像"""
    from config import Config
    import requests
    import json
    
    if not Config.VOLC_SEEDREAM_API_KEY:
        # 如果没有配置API Key，返回优化后的提示词
        from app.ai_service import chat_with_model
        system_prompt = """你是一个专业的图像生成助手。用户会描述他们想要生成的图像，你需要：
1. 理解用户的描述
2. 优化和扩展描述，使其更适合图像生成
3. 提供详细的图像生成提示词（prompt）

注意：你只负责生成提示词，不直接生成图像。请用中文回复。"""
        full_message = f"{system_prompt}\n\n用户需求：{prompt}"
        response = chat_with_model(full_message, 'aliyun-qwen-turbo')
        if response['success']:
            return jsonify({
                'success': True,
                'output': response['message'],
                'type': 'prompt_only',
                'message': '提示：由于未配置图像生成API密钥，仅提供优化后的提示词。请配置VOLC_SEEDREAM_API_KEY以使用图像生成功能。'
            })
        else:
            return jsonify({
                'success': False,
                'message': '图像生成提示词优化失败'
            })
    
    try:
        # 火山引擎Seedream API调用
        # 根据火山引擎API文档，这里使用标准的图像生成接口
        api_url = "https://ark.cn-beijing.volces.com/api/v3/images/generations"
        headers = {
            "Authorization": f"Bearer {Config.VOLC_SEEDREAM_API_KEY}",
            "Content-Type": "application/json"
        }
        
        # 尝试多个可能的模型名称和端点（根据错误信息，模型名称可能不正确）
        api_endpoints = [
            "https://ark.cn-beijing.volces.com/api/v3/images/generations",
            "https://ark.volces.com/api/v3/images/generations",
            "https://open.volcengine.com/api/v3/images/generations"
        ]
        
        # 尝试多个可能的模型名称
        model_names = ["seedream-v1.1", "seedream-v1", "seedream", "seedream-v2", "seedream-pro"]
        
        last_error = None
        success_response = None
        success_model = None
        
        # 尝试所有可能的端点和模型组合
        for api_url in api_endpoints:
            for model_name in model_names:
                try:
                    data = {
                        "model": model_name,
                        "prompt": prompt,
                        "n": 1,
                        "size": "1024x1024",
                        "response_format": "url"
                    }
                    
                    response = requests.post(api_url, headers=headers, json=data, timeout=30)
                    
                    if response.status_code == 200:
                        result = response.json()
                        if 'data' in result and len(result['data']) > 0:
                            success_response = response
                            success_model = model_name
                            break
                    else:
                        # 记录错误信息
                        try:
                            error_json = response.json()
                            last_error = error_json.get('error', {}).get('message', response.text)
                        except:
                            last_error = response.text
                except Exception as e:
                    last_error = str(e)
                    continue
            else:
                continue
            break
        
        # 如果所有尝试都失败
        if success_response is None:
            return jsonify({
                'success': False,
                'message': f'图像生成失败：所有API端点和模型名称尝试均失败。最后错误：{last_error}。\n\n可能的原因：\n1. API密钥无效或未正确配置\n2. 模型名称不正确（已尝试：{", ".join(model_names)}）\n3. API端点不正确\n\n建议：请检查火山引擎控制台，确认正确的模型名称和API端点，或联系火山引擎技术支持。'
            })
        
        # 处理成功的响应
        result = success_response.json()
        if 'data' in result and len(result['data']) > 0:
            image_url = result['data'][0].get('url', '')
            
            # 下载图片并保存到AIGCimages文件夹
            import os
            from datetime import datetime
            from flask_login import current_user
            
            # 创建AIGCimages文件夹（与run.py同目录）
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            images_dir = os.path.join(base_dir, 'AIGCimages')
            os.makedirs(images_dir, exist_ok=True)
            
            # 下载图片
            img_response = requests.get(image_url, timeout=30)
            if img_response.status_code == 200:
                # 生成文件名：用户ID_时间戳_随机数.png
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                user_id = current_user.id if current_user.is_authenticated else 0
                import random
                random_num = random.randint(1000, 9999)
                filename = f"{user_id}_{timestamp}_{random_num}.png"
                filepath = os.path.join(images_dir, filename)
                
                # 保存图片
                with open(filepath, 'wb') as f:
                    f.write(img_response.content)
                
                # 生成相对URL路径
                relative_url = f"/AIGCimages/{filename}"
                
                # 保存历史记录
                save_ai_usage_history(
                    user_id=current_user.id if current_user.is_authenticated else None,
                    tool_type='image-gen',
                    input_text=prompt,
                    output_text=f'图像生成成功',
                    image_url=relative_url,
                    model_used=success_model
                )
                
                return jsonify({
                    'success': True,
                    'output': f'图像生成成功！\n\n提示词：{prompt}',
                    'image_url': relative_url,
                    'type': 'image',
                    'prompt': prompt
                })
            else:
                # 如果下载失败，仍返回原始URL
                return jsonify({
                    'success': True,
                    'output': f'图像生成成功！\n\n提示词：{prompt}\n\n图像URL：{image_url}',
                    'image_url': image_url,
                    'type': 'image',
                    'prompt': prompt
                })
        else:
            return jsonify({
                'success': False,
                'message': '图像生成失败：API返回数据格式错误'
            })
    except requests.exceptions.Timeout:
        return jsonify({
            'success': False,
            'message': '图像生成超时，请稍后重试'
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'图像生成失败：{str(e)}'
        })


def save_conversation_history(user_id, tool_type, conversation_data, model_used=None, conversation_id=None):
    """保存对话历史记录（支持多轮对话）"""
    try:
        from app.models import AIConversationHistory
        from app import db
        import json
        
        # 生成标题（使用第一轮对话的前30个字符）
        title = None
        if conversation_data and len(conversation_data) > 0:
            first_user_msg = next((msg['content'] for msg in conversation_data if msg['role'] == 'user'), '')
            if first_user_msg:
                title = first_user_msg[:30] + ('...' if len(first_user_msg) > 30 else '')
        
        if conversation_id:
            # 更新现有记录
            history = AIConversationHistory.query.filter_by(
                id=conversation_id,
                user_id=user_id
            ).first()
            if history:
                history.conversation_data = json.dumps(conversation_data, ensure_ascii=False)
                history.updated_at = datetime.utcnow()
                if title:
                    history.title = title
                db.session.commit()
                return history.id
        else:
            # 创建新记录
            history = AIConversationHistory(
                user_id=user_id,
                tool_type=tool_type,
                title=title or f'{tool_type}对话',
                conversation_data=json.dumps(conversation_data, ensure_ascii=False),
                model_used=model_used
            )
            db.session.add(history)
            db.session.commit()
            return history.id
    except Exception as e:
        print(f"保存对话历史失败: {str(e)}")
        return None


def get_or_create_conversation_id(user_id, tool_type, conversation_id=None):
    """获取或创建对话ID"""
    if conversation_id:
        return conversation_id
    try:
        from app.models import AIConversationHistory
        from app import db
        # 获取用户最新的该工具类型的对话记录
        latest = AIConversationHistory.query.filter_by(
            user_id=user_id,
            tool_type=tool_type
        ).order_by(AIConversationHistory.created_at.desc()).first()
        return latest.id if latest else None
    except:
        return None


def save_ai_usage_history(user_id, tool_type, input_text, output_text=None, image_url=None, file_path=None, model_used=None):
    """保存工具使用历史记录到数据库"""
    if not user_id:
        return
    
    try:
        from app.models import AIUsageHistory, db
        from app import create_app
        from config import Config
        
        # 检查是否配置了数据库
        if hasattr(Config, 'SQLALCHEMY_DATABASE_URI') and Config.SQLALCHEMY_DATABASE_URI:
            app = create_app()
            with app.app_context():
                history = AIUsageHistory(
                    user_id=user_id,
                    tool_type=tool_type,
                    input_text=input_text[:2000] if input_text else None,
                    output_text=output_text[:5000] if output_text else None,
                    image_url=image_url,
                    file_path=file_path,
                    model_used=model_used
                )
                db.session.add(history)
                db.session.commit()
        else:
            # 如果数据库未配置，使用mock数据存储
            from app.mock_data import MOCK_AI_USAGE_HISTORY
            from datetime import datetime
            history_item = {
                'id': len(MOCK_AI_USAGE_HISTORY) + 1,
                'user_id': user_id,
                'tool_type': tool_type,
                'input_text': input_text[:2000] if input_text else None,
                'output_text': output_text[:5000] if output_text else None,
                'image_url': image_url,
                'file_path': file_path,
                'model_used': model_used,
                'created_at': datetime.now()
            }
            MOCK_AI_USAGE_HISTORY.append(history_item)
    except Exception as e:
        # 如果数据库未配置，使用mock数据存储
        try:
            from app.mock_data import MOCK_AI_USAGE_HISTORY
            from datetime import datetime
            history_item = {
                'id': len(MOCK_AI_USAGE_HISTORY) + 1,
                'user_id': user_id,
                'tool_type': tool_type,
                'input_text': input_text[:2000] if input_text else None,
                'output_text': output_text[:5000] if output_text else None,
                'image_url': image_url,
                'file_path': file_path,
                'model_used': model_used,
                'created_at': datetime.now()
            }
            MOCK_AI_USAGE_HISTORY.append(history_item)
        except:
            pass


def detect_code_language(text):
    """检测代码语言类型"""
    import re
    
    # 检查代码块中的语言标识
    code_block_pattern = r'```(\w+)?\n'
    matches = re.findall(code_block_pattern, text)
    if matches:
        lang = matches[0].lower() if matches[0] else ''
        if lang:
            return lang
    
    # 检查代码块结束标记前的语言
    code_block_pattern2 = r'```(\w+)?\s*\n'
    match = re.search(code_block_pattern2, text)
    if match and match.group(1):
        return match.group(1).lower()
    
    # 根据关键词和代码特征判断
    text_lower = text.lower()
    
    # Python特征
    if any(keyword in text for keyword in ['def ', 'import ', 'from ', 'print(', '__init__', 'if __name__']):
        return 'python'
    # JavaScript特征
    elif any(keyword in text for keyword in ['function ', 'const ', 'let ', 'var ', '=>', 'console.log']):
        return 'javascript'
    # Java特征
    elif 'public class' in text or 'public static void main' in text:
        return 'java'
    # C/C++特征
    elif '#include' in text or 'using namespace' in text:
        return 'cpp'
    # HTML特征
    elif '<html' in text_lower or '<div' in text_lower or '<body' in text_lower:
        return 'html'
    # CSS特征
    elif re.search(r'\{[^}]*:[^}]*\}', text) and ('color:' in text_lower or 'margin:' in text_lower):
        return 'css'
    # SQL特征
    elif any(keyword in text_lower for keyword in ['select ', 'from ', 'where ', 'insert into', 'create table']):
        return 'sql'
    # TypeScript特征
    elif 'interface ' in text or 'type ' in text or ': string' in text or ': number' in text:
        return 'typescript'
    # Go特征
    elif 'package ' in text or 'func ' in text or 'import (' in text:
        return 'go'
    # Rust特征
    elif 'fn ' in text and 'let ' in text:
        return 'rust'
    
    return 'txt'


def chat_with_model_multi_turn(messages, model='aliyun-qwen-turbo'):
    """支持多轮对话的模型调用"""
    from app.ai_service import chat_with_model_multi_turn as ai_multi_turn
    
    # 尝试使用AI服务的多轮对话函数
    try:
        return ai_multi_turn(messages, model)
    except:
        # 如果AI服务没有多轮对话函数，回退到单轮模式
        from app.ai_service import chat_with_model
        
        # 将对话历史转换为上下文文本
        # 保留最近的对话历史（最多10轮）
        recent_messages = messages[-10:] if len(messages) > 10 else messages
        
        conversation_text = ""
        for msg in recent_messages[1:]:  # 跳过system消息
            role = "用户" if msg['role'] == 'user' else "助手"
            conversation_text += f"{role}：{msg['content']}\n\n"
        
        # 构建包含上下文的完整消息
        system_prompt = recent_messages[0]['content'] if recent_messages and recent_messages[0].get('role') == 'system' else "你是一个专业的助手。"
        full_message = f"{system_prompt}\n\n以下是对话历史：\n{conversation_text}\n请根据对话历史回答用户的最新问题。"
        
        return chat_with_model(full_message, model)




@api_bp.route('/playground/code/download', methods=['POST'])
@login_required
def download_code():
    """下载代码文件"""
    from flask import send_file
    from io import BytesIO
    
    data = request.get_json()
    code_content = data.get('code', '')
    extension = data.get('extension', 'txt')
    
    if not code_content:
        return jsonify({
            'success': False,
            'message': '代码内容不能为空'
        })
    
    try:
        # 提取代码块中的代码（如果有代码块标记）
        import re
        
        # 尝试提取代码块
        code_block_pattern = r'```(?:\w+)?\s*\n(.*?)```'
        matches = re.findall(code_block_pattern, code_content, re.DOTALL)
        if matches:
            # 使用第一个代码块
            code_content = matches[0].strip()
        else:
            # 如果没有代码块标记，尝试提取可能的代码部分
            # 查找包含代码特征的行
            lines = code_content.split('\n')
            code_lines = []
            in_code = False
            for line in lines:
                # 检测代码特征
                if any(keyword in line for keyword in ['def ', 'function', 'class ', 'import ', 'const ', 'let ', 'var ', 'public ', 'private ']):
                    in_code = True
                if in_code:
                    code_lines.append(line)
                # 如果遇到空行且已有代码，可能代码结束
                if in_code and not line.strip() and len(code_lines) > 5:
                    break
            
            if code_lines:
                code_content = '\n'.join(code_lines).strip()
        
        # 创建文件
        code_io = BytesIO()
        code_io.write(code_content.encode('utf-8'))
        code_io.seek(0)
        
        # 确定MIME类型
        mime_types = {
            'py': 'text/x-python',
            'js': 'text/javascript',
            'ts': 'text/typescript',
            'java': 'text/x-java-source',
            'cpp': 'text/x-c++',
            'c': 'text/x-c',
            'html': 'text/html',
            'css': 'text/css',
            'sql': 'text/x-sql',
            'sh': 'text/x-shellscript',
            'json': 'application/json',
            'xml': 'application/xml',
            'yml': 'text/yaml',
            'md': 'text/markdown',
            'txt': 'text/plain'
        }
        
        mime_type = mime_types.get(extension, 'text/plain')
        
        # 保存代码文件到本地
        from flask_login import current_user
        if current_user.is_authenticated:
            import os
            from datetime import datetime
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            code_dir = os.path.join(base_dir, 'AIGCimages')
            os.makedirs(code_dir, exist_ok=True)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            user_id = current_user.id
            import random
            random_num = random.randint(1000, 9999)
            code_filename = f"{user_id}_{timestamp}_{random_num}.{extension}"
            code_filepath = os.path.join(code_dir, code_filename)
            with open(code_filepath, 'w', encoding='utf-8') as f:
                f.write(code_content)
            
            # 保存历史记录
            save_ai_usage_history(
                user_id=current_user.id,
                tool_type='programming',
                input_text=data.get('input', '')[:2000] if data.get('input') else None,
                output_text=code_content[:5000],
                file_path=f"/AIGCimages/{code_filename}",
                model_used=data.get('model', 'unknown')
            )
        
        return send_file(
            code_io,
            mimetype=mime_type,
            as_attachment=True,
            download_name=f'code.{extension}'
        )
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'生成代码文件失败: {str(e)}'
        })


@api_bp.route('/playground/conversation/history', methods=['GET'])
@login_required
def get_conversation_history():
    """获取用户的对话历史记录列表"""
    try:
        from app.models import AIConversationHistory
        from flask_login import current_user
        import json
        
        tool_type = request.args.get('tool_type')  # 可选：按工具类型筛选
        page = request.args.get('page', 1, type=int)
        per_page = request.args.get('per_page', 20, type=int)
        
        query = AIConversationHistory.query.filter_by(user_id=current_user.id)
        if tool_type:
            query = query.filter_by(tool_type=tool_type)
        
        pagination = query.order_by(AIConversationHistory.updated_at.desc()).paginate(
            page=page, per_page=per_page, error_out=False
        )
        
        history_list = []
        for history in pagination.items:
            conversation_data = json.loads(history.conversation_data) if history.conversation_data else []
            history_list.append({
                'id': history.id,
                'tool_type': history.tool_type,
                'title': history.title,
                'model_used': history.model_used,
                'created_at': history.created_at.strftime('%Y-%m-%d %H:%M:%S'),
                'updated_at': history.updated_at.strftime('%Y-%m-%d %H:%M:%S'),
                'message_count': len(conversation_data)
            })
        
        return jsonify({
            'success': True,
            'history': history_list,
            'total': pagination.total,
            'page': page,
            'per_page': per_page,
            'pages': pagination.pages
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'获取历史记录失败: {str(e)}'
        })


@api_bp.route('/playground/conversation/<int:conversation_id>', methods=['GET'])
@login_required
def get_conversation_detail(conversation_id):
    """获取单个对话历史记录的详情"""
    try:
        from app.models import AIConversationHistory
        from flask_login import current_user
        import json
        
        history = AIConversationHistory.query.filter_by(
            id=conversation_id,
            user_id=current_user.id
        ).first()
        
        if not history:
            return jsonify({
                'success': False,
                'message': '对话记录不存在'
            })
        
        conversation_data = json.loads(history.conversation_data) if history.conversation_data else []
        
        return jsonify({
            'success': True,
            'id': history.id,
            'tool_type': history.tool_type,
            'title': history.title,
            'model_used': history.model_used,
            'conversation': conversation_data,
            'created_at': history.created_at.strftime('%Y-%m-%d %H:%M:%S'),
            'updated_at': history.updated_at.strftime('%Y-%m-%d %H:%M:%S')
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'获取对话详情失败: {str(e)}'
        })


@api_bp.route('/playground/conversation/<int:conversation_id>', methods=['DELETE'])
@login_required
def delete_conversation(conversation_id):
    """删除对话历史记录"""
    try:
        from app.models import AIConversationHistory
        from app import db
        from flask_login import current_user
        
        history = AIConversationHistory.query.filter_by(
            id=conversation_id,
            user_id=current_user.id
        ).first()
        
        if not history:
            return jsonify({
                'success': False,
                'message': '对话记录不存在'
            })
        
        db.session.delete(history)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '删除成功'
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'删除失败: {str(e)}'
        })


@api_bp.route('/playground/ppt/download', methods=['POST'])
@login_required
def download_ppt():
    """生成并下载PPT文件"""
    from pptx import Presentation
    from io import BytesIO
    from flask import send_file
    import os
    
    data = request.get_json()
    ppt_content = data.get('content', '')
    
    if not ppt_content:
        return jsonify({
            'success': False,
            'message': 'PPT内容不能为空'
        })
    
    try:
        # 创建新的PPT，不使用模板
        prs = Presentation()
        
        # 解析Markdown格式的内容，自动检测和调整页面
        slides_data = []
        slides = ppt_content.split('#')
        
        for slide_content in slides:
            slide_content = slide_content.strip()
            if not slide_content:
                continue
            
            lines = slide_content.split('\n')
            title = lines[0].strip()
            if not title:
                continue
            
            # 清理标题：移除"幻灯片X："这样的前缀
            import re
            title = re.sub(r'^幻灯片\d+[：:]\s*', '', title)
            title = re.sub(r'^Slide\s+\d+[：:]\s*', '', title, flags=re.IGNORECASE)
            
            # 提取要点
            bullet_points = []
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                if line.startswith('-') or line.startswith('•') or line.startswith('*'):
                    bullet_points.append(line.lstrip('-•*').strip())
                elif line:
                    bullet_points.append(line)
            
            # 自动检测内容是否适合单页，如果内容过多则自动分页
            MAX_POINTS_PER_SLIDE = 6  # 每页最多6个要点
            MAX_TITLE_LENGTH = 50  # 标题最大长度
            MAX_POINT_LENGTH = 100  # 每个要点最大长度
            
            # 如果标题过长，截断
            if len(title) > MAX_TITLE_LENGTH:
                title = title[:MAX_TITLE_LENGTH] + '...'
            
            # 如果要点过多或单个要点过长，需要分页
            if len(bullet_points) > MAX_POINTS_PER_SLIDE:
                # 分页处理
                total_pages = (len(bullet_points) + MAX_POINTS_PER_SLIDE - 1) // MAX_POINTS_PER_SLIDE
                for i in range(0, len(bullet_points), MAX_POINTS_PER_SLIDE):
                    page_points = bullet_points[i:i + MAX_POINTS_PER_SLIDE]
                    # 截断过长的要点
                    page_points = [p[:MAX_POINT_LENGTH] + '...' if len(p) > MAX_POINT_LENGTH else p for p in page_points]
                    
                    page_title = title
                    if total_pages > 1:
                        page_title = f"{title}（{i//MAX_POINTS_PER_SLIDE + 1}/{total_pages}）"
                    
                    slides_data.append({
                        'title': page_title,
                        'points': page_points
                    })
            else:
                # 单页内容，截断过长的要点
                bullet_points = [p[:MAX_POINT_LENGTH] + '...' if len(p) > MAX_POINT_LENGTH else p for p in bullet_points]
                
                slides_data.append({
                    'title': title,
                    'points': bullet_points
                })
        
        # 如果没有幻灯片，创建一个默认的
        if not slides_data:
            slides_data.append({
                'title': 'PPT内容',
                'points': []
            })
        
        # 创建幻灯片并添加内容
        for i, slide_data in enumerate(slides_data):
            # 选择合适的布局
            # 0: 标题幻灯片, 1: 标题和内容, 5: 仅标题
            if i == 0 and len(slides_data) > 1:
                # 第一张幻灯片使用标题幻灯片（布局0）
                layout_index = 0
            else:
                # 其他幻灯片使用标题和内容布局（布局1）
                layout_index = 1
            
            # 确保布局索引有效
            if layout_index >= len(prs.slide_layouts):
                layout_index = 0
            
            # 创建幻灯片
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            except:
                # 如果布局不存在，使用第一个可用布局
                slide = prs.slides.add_slide(prs.slide_layouts[0])
            
            # 填充标题
            title_added = False
            try:
                # 尝试使用标题占位符
                if hasattr(slide, 'shapes') and slide.shapes.title:
                    title_shape = slide.shapes.title
                    if title_shape.has_text_frame:
                        title_shape.text = slide_data['title']
                        title_added = True
            except:
                pass
            
            # 如果没有标题占位符，尝试查找其他文本框
            if not title_added:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        try:
                            # 检查是否是标题位置（通常在顶部）
                            if shape.top < slide.height * 0.2:  # 在顶部20%区域内
                                shape.text_frame.text = slide_data['title']
                                title_added = True
                                break
                        except:
                            pass
            
            # 填充内容
            if slide_data['points']:
                content_added = False
                
                # 方法1: 尝试使用内容占位符（索引1）
                try:
                    if len(slide.placeholders) > 1:
                        content_placeholder = slide.placeholders[1]
                        if content_placeholder != slide.shapes.title and hasattr(content_placeholder, 'text_frame'):
                            text_frame = content_placeholder.text_frame
                            text_frame.clear()
                            
                            # 设置第一个段落
                            if len(text_frame.paragraphs) > 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            
                            p.text = slide_data['points'][0]
                            p.level = 0
                            
                            # 添加其他要点
                            for point in slide_data['points'][1:]:
                                p = text_frame.add_paragraph()
                                p.text = point
                                p.level = 0
                            
                            content_added = True
                except:
                    pass
                
                # 方法2: 如果没有找到占位符，查找内容区域的文本框
                if not content_added:
                    for shape in slide.shapes:
                        if shape != slide.shapes.title and hasattr(shape, 'text_frame'):
                            text_frame = shape.text_frame
                            if text_frame:
                                try:
                                    # 检查是否是内容位置（不在顶部）
                                    if shape.top > slide.height * 0.15:  # 在顶部15%以下
                                        text_frame.clear()
                                        
                                        # 设置第一个段落
                                        if len(text_frame.paragraphs) > 0:
                                            p = text_frame.paragraphs[0]
                                        else:
                                            p = text_frame.add_paragraph()
                                        
                                        p.text = slide_data['points'][0]
                                        p.level = 0
                                        
                                        # 添加其他要点
                                        for point in slide_data['points'][1:]:
                                            p = text_frame.add_paragraph()
                                            p.text = point
                                            p.level = 0
                                        
                                        content_added = True
                                        break
                                except:
                                    pass
        
        # 保存到内存
        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        
        # 保存PPT文件到本地
        from flask_login import current_user
        if current_user.is_authenticated:
            import os
            from datetime import datetime
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            ppt_dir = os.path.join(base_dir, 'AIGCimages')
            os.makedirs(ppt_dir, exist_ok=True)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            user_id = current_user.id
            import random
            random_num = random.randint(1000, 9999)
            ppt_filename = f"{user_id}_{timestamp}_{random_num}.pptx"
            ppt_filepath = os.path.join(ppt_dir, ppt_filename)
            with open(ppt_filepath, 'wb') as f:
                f.write(ppt_io.getvalue())
            
            # 保存历史记录
            save_ai_usage_history(
                user_id=current_user.id,
                tool_type='ppt',
                input_text=data.get('input', '')[:2000] if data.get('input') else None,
                output_text=ppt_content[:5000],
                file_path=f"/AIGCimages/{ppt_filename}",
                model_used=data.get('model', 'unknown')
            )
        
        return send_file(
            ppt_io,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation.pptx'
        )
    except ImportError:
        return jsonify({
            'success': False,
            'message': 'PPT生成功能需要安装python-pptx库，请运行: pip install python-pptx'
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'生成PPT失败: {str(e)}'
        })


@api_bp.route('/ai-usage-history', methods=['GET'])
@login_required
def get_ai_usage_history():
    """获取工具使用历史记录"""
    from flask_login import current_user
    from datetime import datetime
    
    try:
        from app.models import AIUsageHistory, db
        from app import create_app
        from config import Config
        
        # 检查是否配置了数据库
        if hasattr(Config, 'SQLALCHEMY_DATABASE_URI') and Config.SQLALCHEMY_DATABASE_URI:
            app = create_app()
            with app.app_context():
                histories = AIUsageHistory.query.filter_by(user_id=current_user.id)\
                    .order_by(AIUsageHistory.created_at.desc())\
                    .limit(100).all()
                
                result = []
                for h in histories:
                    result.append({
                        'id': h.id,
                        'tool_type': h.tool_type,
                        'input_text': h.input_text,
                        'output_text': h.output_text[:200] if h.output_text else None,  # 只返回前200字符
                        'image_url': h.image_url,
                        'file_path': h.file_path,
                        'model_used': h.model_used,
                        'created_at': h.created_at.strftime('%Y-%m-%d %H:%M:%S') if h.created_at else None
                    })
                return jsonify({'success': True, 'data': result})
        else:
            # 使用mock数据
            from app.mock_data import MOCK_AI_USAGE_HISTORY
            user_histories = [h for h in MOCK_AI_USAGE_HISTORY if h.get('user_id') == current_user.id]
            user_histories.sort(key=lambda x: x.get('created_at', datetime.now()), reverse=True)
            user_histories = user_histories[:100]  # 限制100条
            
            result = []
            for h in user_histories:
                result.append({
                    'id': h.get('id'),
                    'tool_type': h.get('tool_type'),
                    'input_text': h.get('input_text'),
                    'output_text': h.get('output_text', '')[:200] if h.get('output_text') else None,
                    'image_url': h.get('image_url'),
                    'file_path': h.get('file_path'),
                    'model_used': h.get('model_used'),
                    'created_at': h.get('created_at').strftime('%Y-%m-%d %H:%M:%S') if isinstance(h.get('created_at'), datetime) else str(h.get('created_at', ''))
                })
            return jsonify({'success': True, 'data': result})
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'获取历史记录失败: {str(e)}'
        })
